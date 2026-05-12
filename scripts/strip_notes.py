"""
Strip all speaker notes from WSU-Workshop-May15.pptx.

The presenter has their own notes (WSU-Workshop-Speech.docx) and does not want
the per-slide notes baked into the PowerPoint file. Run this AFTER build_pptx.py.
"""
import os
from pptx import Presentation

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
PPTX = os.path.join(ROOT, 'WSU-Workshop-May15.pptx')

prs = Presentation(PPTX)
for i, slide in enumerate(prs.slides):
    nf = slide.notes_slide.notes_text_frame
    # Clear text by setting it to an empty string
    nf.text = ''

prs.save(PPTX)
print(f"Stripped notes from {len(prs.slides)} slides in {PPTX}")
