"""
Replace the speaker notes in WSU-Workshop-May15.pptx with conversational,
first-person scripts the presenter will read aloud to students.

Each note is what Cagri actually says to the room, not a stage direction.
"""
import os
from pptx import Presentation

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
PPTX = os.path.join(ROOT, 'WSU-Workshop-May15.pptx')

# Notes in deck order. Indexed by slide number (0-based).
NOTES = [

    # 0. Cover
    "Good morning everyone. I'm Cagri Temel, CTO at Hezarfen LLC. "
    "For the next two hours and forty-five minutes you're going to do something most people "
    "don't get to do until they're in industry: you're going to train a real machine learning "
    "model on real NASA data, you're going to break it on purpose, and then you're going to "
    "explain exactly how it makes its decisions. By the end of the session you'll have a working "
    "model on your laptop, you'll know which sensor failed on an aircraft engine, and many of you "
    "will have your first open-source contribution on GitHub. Let's get started.",

    # 1. Who I am
    "Quick about me. I'm CTO of a company called Hezarfen, where we build explainable AI for industries "
    "that have to answer to regulators: banks, insurers, healthcare, aviation. I'm an IEEE Senior Member. "
    "I also maintain a Python package called neural-trees, and that's what we'll use today. "
    "Don't memorize any of this. The only thing I want you to take away from this slide is the bottom line: "
    "in two hours and forty-five minutes, you build, you break, and you explain a real model.",

    # 2. Today's mission
    "For the next few hours we're going to look at one specific question. "
    "How do we make data-driven systems trustworthy in places where getting them wrong is expensive? "
    "Our case study today is a turbofan engine, the kind that powers passenger jets, "
    "that has to predict its own failure before it happens. "
    "By the end of the session you'll understand why this is harder than it looks, "
    "and why being accurate is not enough.",

    # 3. Hook
    "So. A turbofan engine. The question is, when will it fail? "
    "Think about it for a moment. If you knew with certainty, exactly when this engine would fail, "
    "what would you do differently than what airlines do today? "
    "I'm not going to call on anybody. Just hold that question in your head as we go.",

    # 4. Stakes
    "Here's why this is hard. There are three ways to be wrong, and all three are expensive. "
    "One: you predict too late. The engine fails in flight. You lose lives. "
    "Two: you predict too early. You pull a perfectly good engine out of service. "
    "Each engine is about a hundred million dollars. "
    "Three, and this is the one engineers forget: you predict accurately, but you can't explain how. "
    "In that case, the FAA refuses to certify your system. "
    "You can have the best model in the world; if you can't defend it, it doesn't fly. "
    "That's the whole point. Accuracy by itself is not enough. You also need to defend the prediction.",

    # 5. CMAPSS
    "The dataset we'll use is called NASA CMAPSS. It's open, it's free, and it's been the standard "
    "benchmark in predictive maintenance for about fifteen years. A hundred turbofan engines, each one "
    "simulated all the way to failure. Twenty-one sensors per timestep: temperatures, pressures, fan speeds, "
    "fuel flow. Each engine logs hundreds of flight cycles before it dies. "
    "You'll be working with this data on your own laptop in twenty minutes.",

    # 6. RUL
    "The thing we're going to predict is called RUL, Remaining Useful Life. "
    "It's the number of flight cycles left before the engine fails. "
    "So if an engine fails at cycle 192, at cycle 100 its RUL is 92. At cycle 190 it's 2. At failure it's zero. "
    "We cap RUL at 125 because predicting further than that ahead isn't useful. "
    "Past 125 cycles, you're just guessing about the future.",

    # 7. Sensors
    "For every cycle, we get 21 sensor readings. Don't try to memorize them. "
    "But three of them I want you to notice now, because they'll come back later. "
    "Sensor 7, high-pressure compressor outlet pressure. "
    "Sensor 11, static pressure. "
    "Sensor 14, corrected core speed. "
    "These three are the model's favorites. Remember those names. They show up again in Activity 2 "
    "when one of them gets attacked.",

    # 8. LSTM
    "Here's the modern approach. We throw an LSTM at it. A two-layer recurrent neural network. "
    "After training, the LSTM gets an RMSE of 15.49 cycles, R-squared 0.855. That's pretty good. "
    "About fifteen cycles off, on average. So... we ship it to the FAA, right?",

    # 9. Regulator
    "No. "
    "Regulators don't ask, how accurate is your model. They ask, why did your model say that? "
    "And if you can't answer that question in a way a non-engineer can follow, your model doesn't ship. "
    "FAA for aviation. FDA for medical devices. OCC for banking. EU AI Act for anything that touches a European customer. "
    "This isn't theory. This is what actually blocks AI systems from being deployed in regulated industries.",

    # 10. Black box vs glass box
    "So you're stuck between two worlds. On one side, black-box models like LSTMs and transformers. High accuracy, "
    "but you can't explain why. On the other side, glass-box models like classic decision trees or linear regression. "
    "Fully explainable, but lower accuracy. Regulators love glass boxes. Engineers prefer black boxes. "
    "So the question today is: can we have both?",

    # 11. Classic Decision Tree
    "You've all seen a decision tree before. Beautiful structure. Sensor 11 is above the threshold, go left; below, go right. "
    "It's the most explainable model that exists. But classic decision trees often have worse accuracy than neural networks. "
    "Why? Because the splits are hard. A 0.01 difference in a sensor value flips the entire path. "
    "That brittleness is what limits accuracy.",

    # 12. Hard split
    "This is the code for a hard split. You are NOT typing this. It's here for illustration. "
    "If x is bigger than the threshold, return right. Otherwise, return left. "
    "The output is a string. Discrete. No gradient. "
    "And no gradient means no backpropagation. You build the tree node by node, greedily. "
    "You cannot train it end to end like a neural network.",

    # 13. Soft split
    "Now look at this. Same function. One substitution. The if-else becomes a sigmoid. "
    "Now the output is a probability, somewhere between zero and one. Continuous. "
    "And because it's continuous, gradient flows through it. "
    "That single change is what makes everything we do today possible. "
    "A tree you can train with backpropagation, just like a neural network.",

    # 14. Soft Tree Visual
    "Here's what the whole tree looks like. Every internal node is a sigmoid split. "
    "Every leaf has a class distribution: Critical, Caution, Healthy probabilities. "
    "The final prediction is a weighted sum of all the leaves, weighted by the probability of reaching each leaf. "
    "Mathematically, it's smooth. Fully differentiable end to end.",

    # 15. Best of Both
    "And here's the win. From the tree side, you keep path traceability, per-feature thresholds, "
    "human-readable rules. From the network side, you get end-to-end training and the accuracy that comes with it. "
    "You get both. That's why soft decision trees matter for regulated industries.",

    # 16. Three properties
    "Three properties make this approach work for safety-critical AI. "
    "One, explainability. Every prediction comes with the path that produced it. "
    "Two, noise robustness. Soft splits don't break when a sensor wobbles a little. "
    "Three, sensor failure tolerance. When you train with channel-level dropout, the model learns to operate "
    "even when some sensors are missing. "
    "You'll see all three in the next ninety minutes. Not on a slide. On your own laptop.",

    # 17. Paper numbers
    "These are the numbers from my paper. On clean CMAPSS data, an LSTM gets RMSE 15.49 cycles. "
    "The Temporal Neural Tree, the variant we use, gets 15.78. Essentially equal. "
    "Now look at the next column. When thirty percent of the sensors are missing — which happens in real aircraft, "
    "sensors break — the LSTM degrades by eighty-nine percent. The Temporal Neural Tree degrades by only seventeen percent. "
    "Five times more robust. And it's fully explainable. That's the headline result.",

    # 18. neural-trees package
    "The tool you'll use is called neural-trees. I wrote it. It's on PyPI, MIT licensed, open source. "
    "The API follows scikit-learn, so if you've used scikit-learn before you already know how to use this: "
    "fit, predict, predict_proba, score. The backend is PyTorch. You don't need a GPU; CPU is fine for everything we do today. "
    "And at the end of the workshop, you'll get a chance to contribute back to it.",

    # 19. Activity 1 intro
    "Okay. We're moving from lecture to hands-on now. The next fifty-five minutes are yours. "
    "You're going to train your own neural tree on the NASA data we just talked about. "
    "I'm going to walk around. If you get stuck, raise a hand, I'll come help. "
    "There are eleven steps. Each step is on a slide. You paste the code, run it, and check that the output matches mine. "
    "Ready? Let's go.",

    # 20. QR
    "Pull out your phones and scan this QR code. It opens our workshop landing page. "
    "Tap Activity 1. That opens a Google Colab notebook on your laptop. "
    "Sign in with your Gmail, then click Copy to Drive so you have your own editable copy. "
    "Run the first cell to make sure everything works. "
    "I'll give you two minutes. Raise your hand if you get stuck.",

    # 21. Overview
    "Quick overview of the next fifty-five minutes. Six main stages. "
    "We load CMAPSS and look at it. Bin the labels into three classes. Train the Soft Decision Tree, which takes about a minute. "
    "Evaluate it with a confusion matrix. And then — this is the step that matters most — we traverse one prediction. "
    "We literally read out the rule the model used. "
    "If you remember one thing from this entire workshop, it should be that step. Let's start.",

    # 22. Step 1 — Imports
    "First step, imports. Paste this code into a fresh cell. Hit Shift+Enter. "
    "You should see All set. "
    "If you get an IndentationError, that means one of your lines has a stray space at the beginning. "
    "Copy-paste does that sometimes. Click into the cell, hit Command-A to select everything, then Shift+Tab to clear indentation, "
    "and run again.",

    # 23. Step 2 — Load data
    "Now we load the data. CMAPSS comes as a plain text file with no headers, so we tell pandas what to call the columns. "
    "After this runs, you should see two things. The shape, twenty thousand six hundred and thirty-one rows by twenty-six columns. "
    "And the number of engines, one hundred. The first five rows show you engine 1 at cycle 1 through cycle 5. "
    "Those numbers — temperatures, pressures, fan speeds — those are exactly what a real flight data recorder logs.",

    # 24. Step 3 — Compute RUL
    "Step 3 computes the label we're going to predict. For each row, RUL equals the engine's max cycle minus the current cycle. "
    "We cap it at 125. "
    "After this runs, you'll see five rows of engine 1. All five show RUL equal to 125, because at the start of the engine's life, "
    "the real RUL is way above 125, but we capped it. That's expected.",

    # 25. Step 4 — Plot engine 1
    "Now we visualize. We're going to plot four sensors for engine 1 across its entire one hundred and ninety-two cycle life. "
    "After this runs, you should see a two-by-two plot. Look at it closely. "
    "Sensor 2 — temperature — goes up over time. Sensor 11 — pressure — goes up. Sensor 14 — corrected core speed — goes up. "
    "But sensor 7 — pressure at the HPC outlet — goes down. "
    "That's the engine wearing out. Temperatures rising, compressor pressure dropping. "
    "The data is telling us the story of failure. Now we have to teach a model to see it.",

    # 26. Step 5 — Bin RUL
    "We're going to turn this into a classification problem. RUL is a number, but in real aviation operations, no one says "
    "'this engine has forty-seven cycles left.' They say 'this engine is healthy,' or 'caution,' or 'critical.' "
    "So we bin RUL into three classes. Below 30 is Critical. Thirty to eighty is Caution. Above eighty is Healthy. "
    "When you run this, you should see roughly fourteen percent Critical, twenty-four percent Caution, sixty-one percent Healthy. "
    "The classes are imbalanced. That's reality. Most of the time the engine is fine.",

    # 27. Step 6 — Feature matrix
    "Now we drop the six sensors that are constant. They carry no information. That leaves fifteen informative sensors. "
    "Then we standardize. Every column gets zero mean and unit variance. "
    "Soft decision trees don't strictly require standardization, but the sigmoid splits behave much better when the input is on a consistent scale. "
    "After this runs, you should see X shape twenty thousand by fifteen, y shape twenty thousand.",

    # 28. Step 7 — Split
    "Standard train-test split. Eighty percent training, twenty percent test. "
    "Two things matter. Stratify equals y, which preserves the class ratio in both sets. "
    "And random state equals forty-two, so everyone in this room gets the exact same numbers. "
    "If yours don't match, you typed something wrong. "
    "After this you should see sixteen thousand five hundred and four rows for train, four thousand one hundred and twenty-seven for test.",

    # 29. Step 8 — Train
    "Okay, this is the slow one. About thirty to sixty seconds on CPU. "
    "While we wait, let me tell you what's happening. PyTorch is doing backpropagation on sixteen thousand rows for thirty epochs. "
    "Every epoch, the loss should come down a little. When it finishes, you should see test accuracy around 0.84. "
    "If yours lands anywhere between 0.82 and 0.86, that's normal — randomness in training varies a little across machines. "
    "While we wait — quick question. What do you think an LSTM would get on the same problem? "
    "Right. About the same. The point is not accuracy. The point is what we can do with the model next.",

    # 30. Step 9 — Confusion matrix
    "Now we look at where the model is right and where it's wrong. "
    "You should see a confusion matrix, three by three, color-coded. "
    "Look at the bottom-left cell. It should say zero. That means: out of every Healthy engine in the test set, "
    "the model wrongly flagged zero of them as Critical. We never raised a false alarm. "
    "Now look at the top row. Out of six hundred Critical engines, we caught about four hundred and ninety. We missed some, "
    "but only five of those misses went all the way to Healthy. "
    "In aviation language: we don't miss many dangerous engines, and we don't waste healthy ones. That's a useful model.",

    # 31. Step 10 — Split weights
    "Now we peek inside the model. Soft decision trees have fifteen internal nodes — one root, then it branches out. "
    "For each internal node, we ask: which sensor does this node lean on the most? "
    "When you run this, you'll see fifteen nodes and their dominant sensors. "
    "Notice some sensors show up multiple times. Sensor 14 and sensor 6, in particular. "
    "Those are the pillars of the model. If those sensors go bad, the model is in trouble. "
    "Remember sensor 14. It comes back in Activity 2.",

    # 32. Step 11 — Traverse prediction
    "This is the step that matters most. We pick one specific test sample, sample 17, and we ask the model: "
    "what did you predict, with what confidence, and which sensor did you look at first? "
    "When you run this, you should see four lines. Read them out loud to yourself. "
    "The model predicted Healthy with eighty-five percent confidence. Zero percent Critical. "
    "The root node leaned on sensor 7, HPC outlet pressure. "
    "Now imagine you have to explain this prediction to a regulator. With this output, you can say: "
    "'My model classified this engine as Healthy because the static pressure at the high-pressure compressor outlet "
    "was above its decision threshold.' "
    "An LSTM cannot give you this answer. That's why we did all of this.",

    # 33. Checkpoint
    "Okay. Take a breath. By now you should have three things on your screen: a trained Soft Decision Tree, "
    "test accuracy above eighty percent, and a printed decision path for engine 17. "
    "If you're missing one of those, find me during the break and we'll catch up. "
    "Ten minutes break. Get coffee, stretch. When you come back at ten thirty-five, we run the team competition.",

    # 34. Activity 2 Intro
    "Okay, we're back. Activity 2 is different. This is a competition. You'll work in teams. "
    "The first team to identify all three attacks correctly gets bragging rights for the day. "
    "Here's the setup. One of engine 17's sensors is reporting bad data. "
    "You don't know which sensor. You don't know what kind of attack. Your job is to find both.",

    # 35. Scenario
    "Imagine you're the data science team at an airline. Maintenance ops just sent you a message. "
    "Something is wrong with engine 17. The pattern looks like a sensor drift, a frozen sensor, or noise spikes, "
    "but we can't tell which channel is at fault. Localize the faulty sensor before the next flight. "
    "You have two tools. Your Soft Decision Tree from Activity 1, which is explainable. "
    "And we'll add a RandomForest baseline, which is more accurate but opaque. "
    "Use both. See which one is actually useful.",

    # 36. Teams
    "Okay. Form groups of three to four people. Pick the people next to you. Thirty seconds. "
    "Use one notebook per team. One person types, the others think and argue. "
    "Once we start the clock you have seventeen minutes for analysis. "
    "First team to identify all three correctly — sensor and attack type for each file — wins.",

    # 37. Three attacks
    "Three files. Attack A is a drift: a sensor with a constant offset added. "
    "Attack B is stuck-at: a sensor frozen at a single value. "
    "Attack C is Gaussian noise: a sensor with random noise injected. "
    "All three files are copies of clean engine 17 data, but in each one, exactly one sensor channel has been manipulated. "
    "Find which one.",

    # 38. A2 Step 1 — RF baseline
    "First, we add a RandomForest baseline. Same data, same labels. "
    "When you run this, you should see two accuracies: SoftTree around 0.841, RandomForest around 0.843. "
    "Essentially identical on clean data. The whole point of this activity is showing what happens when we move OFF clean data.",

    # 39. A2 Step 2 — Clean baseline
    "Now we load the clean engine 17 file. This is the reference — the un-attacked version. "
    "We run both models on it and store the predictions. After this you should see two hundred and seventy-six cycles, "
    "and the two models agree on ninety-five percent of them. "
    "That's our baseline. When reality is normal, the two models say the same thing.",

    # 40. A2 Step 3 — Score attacks
    "Now we run both models on each of the three attack files and compare against the clean baseline. "
    "Read your output carefully. "
    "Attack A: tree changed eighteen cycles, RandomForest changed fourteen. "
    "Attack B: tree changed thirty-three, RandomForest changed eight. "
    "Attack C — and this is the one — tree changed eight, RandomForest changed only one. "
    "The RandomForest barely noticed Attack C. It says 'everything looks fine' when in reality a sensor is broken. "
    "The Soft Tree caught eight cycles. The explainable model is also the more sensitive one to subtle manipulations. "
    "That's not a coincidence.",

    # 41. A2 Step 4 — Attack A
    "Now we visualize Attack A. We plot each sensor twice. Clean in blue, attack in orange. "
    "Fourteen of the fifteen panels will show only one line, because the blue and orange overlap exactly. "
    "But ONE panel will show two separate lines. That's your manipulated sensor. "
    "Look carefully. Which sensor is it? "
    "If you see two parallel lines — same shape, but shifted — that's drift. A constant offset. "
    "The answer is sensor 11.",

    # 42. A2 Step 5 — Attack B
    "Same plot, same logic, for Attack B. Look for the sensor where blue and orange diverge. "
    "This time the orange line will be completely flat. Frozen at one value, while the blue line keeps evolving. "
    "That's stuck-at. "
    "Which sensor? Sensor 14. "
    "Now think about this. Sensor 14 was a pillar of our model in Activity 1, Step 10. "
    "Whoever attacked engine 17 went after exactly what the model relies on the most. "
    "That's not a coincidence. That's how adversarial attacks work.",

    # 43. A2 Step 6 — Attack C
    "Attack C is the subtlest one. Look at all fifteen panels. "
    "In one of them, the orange line follows the same trend as blue, but it's noticeably more jagged. "
    "More wobbly. Same average, much higher variance. That's Gaussian noise. "
    "Which sensor? Sensor 9. "
    "This one is the hardest to catch visually because the overall trend is preserved. Only the noise level changes. "
    "That's why the RandomForest missed it. The trend was right, so it kept saying 'fine.'",

    # 44. A2 Step 7 — Quantify
    "Last step. Let's quantify what we just saw with our eyes. "
    "For each attack file, we compute the mean absolute difference between clean and attack for every sensor. "
    "Then we print the top three sensors that differ. "
    "After this runs, you should see: Attack A, sensor 11 first, the rest zero. "
    "Attack B, sensor 14 first with a huge number, nineteen point eight nine, the rest zero. "
    "Attack C, sensor 9 first, six point four two, the rest zero. "
    "Notice that for each attack, only ONE sensor differs at all. The other fourteen are perfectly identical. "
    "Attackers don't manipulate everything. They go after one sensor and try to stay hidden. "
    "Our model caught them anyway.",

    # 45. A2 Clock
    "Okay, seventeen minutes on the clock starting now. Open activity 2 student notebook on your laptop. "
    "Discuss with your team. I'm walking around. If you have questions, raise a hand. "
    "I'll give hints, but I won't give you the answer. Go.",

    # 46. A2 Answers
    "Okay, time's up. Here are the answers. "
    "Attack A is sensor 11, drift. Attack B is sensor 14, stuck-at. Attack C is sensor 9, Gaussian noise. "
    "Show of hands — which team got all three? "
    "Round of applause for them.",

    # 47. What just happened
    "Here's what we just demonstrated. The RandomForest changed its prediction under attack, "
    "but it couldn't tell you which sensor caused the change. "
    "The Soft Decision Tree's split weights shifted in a sensor-specific way. That shift IS the explanation. "
    "It's a fingerprint of which sensor the model is leaning on. "
    "And here's the lesson. Explainability isn't just for regulatory compliance. "
    "It's also a debugging tool. The same property that lets you defend a model to the FAA "
    "also helps your maintenance engineers find the broken hardware. That's the takeaway.",

    # 48. Sprint Intro
    "Okay, last fifteen minutes. This part is different from everything else. "
    "You're going to make a real open-source contribution to a real Python library that's on PyPI. "
    "Not a workshop exercise. A real pull request that I will review and merge live, in front of you, right now. "
    "Your name on it. Forever.",

    # 49. Sprint mechanics
    "Here's how it works. Open github.com/cgrtml/neural-trees. Go to the Issues tab. "
    "Filter by good first issue. Pick one — they're all ten to twenty minute scope. "
    "Comment 'I'm taking this' so two people don't pick the same one. "
    "Fork the repo, make the change in your fork, open a pull request. "
    "I'll review it from my laptop right here on stage. If it's correct, I merge it on the spot. "
    "There are about twenty issues waiting. Pick what fits your level.",

    # 50. Sprint why
    "Why is this worth fifteen minutes of your day? Three reasons. "
    "One. Your contribution shows up on your GitHub profile. That's visible to recruiters, internship coordinators, grad school admissions. "
    "It's a credential that's much harder to fake than a tutorial certificate. "
    "Two. You become part of a Python package that other people install with pip. Your code runs on their machines. "
    "Three. Every merged PR will be credited in the neural-trees README under a section called WSU Workshop Contributors. "
    "That section will exist after today, with your names on it. "
    "Open the repo. Pick an issue. Go.",

    # 51. EU AI Act
    "Quick context for why all of this matters in the real world. "
    "The EU AI Act took effect in 2025. If your AI system serves any European customer, three articles apply. "
    "Article 13: transparency. Article 14: human oversight. Article 15: robustness against adversarial inputs and faults. "
    "The exact things we did today. "
    "And the US is moving in the same direction. SR 11-7 for banking. FDA for medical devices. FAA for aviation. "
    "The era of just-ship-the-LSTM is ending. The era of show-your-work is starting.",

    # 52. Wrap-up
    "Three things to take with you. "
    "One: in safety-critical AI, explainability is not optional. It's a regulatory requirement. "
    "Two: Soft Decision Trees give you both worlds. Neural-network-level accuracy with tree-level interpretability. "
    "Three: you actually did this. Not me. Not a video. You. "
    "You trained a model, you localized a sensor fault, and many of you just shipped your first open-source contribution.",

    # 53. Career bridge
    "Where does this work go next? "
    "Banks, insurers, healthcare, aerospace. Every regulated industry needs explainable AI right now. "
    "My company Hezarfen builds compliance-ready AI for US mid-market banks under SR 11-7 and the EU AI Act. "
    "If anything you saw today interests you — a thesis topic, an open-source contribution, internship questions, "
    "or just a conversation about where to go after graduation — reach out. "
    "My email and LinkedIn are on the next slide.",

    # 54. Links
    "Here are the three things worth bookmarking. "
    "My GitHub, where neural-trees lives. My LinkedIn — happy to connect with any of you. "
    "And the workshop repo on GitHub. Everything we did today is there: the data, the notebooks, the slides. "
    "Take a photo of this slide if you want.",

    # 55. Thank you
    "Thank you. Special thanks to Dr. Sergey Lapin and to Jeremy for the invitation. "
    "Without them this wouldn't have happened. "
    "I'll stay for five minutes after the official close. If you have questions, come up to the stage. "
    "If you want to follow up later, you have my email.",
]

assert len(NOTES) == 56, f"Expected 56 notes, got {len(NOTES)}"

prs = Presentation(PPTX)
assert len(prs.slides) == 56, f"Expected 56 slides, got {len(prs.slides)}"

for i, slide in enumerate(prs.slides):
    nf = slide.notes_slide.notes_text_frame
    # Clear existing notes
    for p in list(nf.paragraphs):
        p.clear()
    nf.text = NOTES[i]

prs.save(PPTX)
print(f"Patched {len(NOTES)} speaker notes in {PPTX}")
