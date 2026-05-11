"""
Build a PowerPoint version of the WSU workshop deck with speaker notes for each slide.

Output: /Users/mac/Projects/wsu-workshop-may15/WSU-Workshop-May15.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
IMG = os.path.join(ROOT, 'slides', 'img')
OUT_PATH = os.path.join(ROOT, 'WSU-Workshop-May15.pptx')

# ===== color palette =====
BG = RGBColor(0x0A, 0x0A, 0x0A)
BG_SECTION = RGBColor(0x10, 0x1A, 0x10)
BG_WARN = RGBColor(0x1A, 0x0A, 0x0A)
ACCENT = RGBColor(0x4F, 0xC3, 0xF7)
ACCENT_DARK = RGBColor(0x00, 0x6A, 0xFF)
TEXT = RGBColor(0xEE, 0xEE, 0xEE)
MUTED = RGBColor(0xBB, 0xBB, 0xBB)
DIM = RGBColor(0x88, 0x88, 0x88)
GOOD = RGBColor(0x81, 0xC7, 0x84)
WARN = RGBColor(0xFF, 0xB7, 0x4D)
BAD = RGBColor(0xEF, 0x53, 0x50)
CODE_BG = RGBColor(0x14, 0x14, 0x14)
OUT_BG = RGBColor(0x0E, 0x1D, 0x0E)
TAG_BG = RGBColor(0x4F, 0xC3, 0xF7)
TAG_BG_A2 = RGBColor(0xBA, 0x68, 0xC8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

BLANK = 6  # blank layout index in default template


def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height, *,
             size=18, bold=False, italic=False, color=TEXT, align=PP_ALIGN.LEFT,
             font='Helvetica Neue', anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_multi(slide, runs, left, top, width, height, *,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    """runs = list of (text, dict_of_options) where each run is on a new paragraph."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.vertical_anchor = anchor
    for i, (text, opts) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get('align', align)
        p.line_spacing = line_spacing
        if 'space_after' in opts:
            p.space_after = Pt(opts['space_after'])
        run = p.add_run()
        run.text = text
        run.font.name = opts.get('font', 'Helvetica Neue')
        run.font.size = Pt(opts.get('size', 16))
        run.font.bold = opts.get('bold', False)
        run.font.italic = opts.get('italic', False)
        run.font.color.rgb = opts.get('color', TEXT)
    return tb


def add_tag(slide, text, left, top, *, bg=TAG_BG):
    """Small accent label, e.g. STEP 1 / 11."""
    width = Inches(2.4)
    height = Inches(0.42)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    shape.adjustments[0] = 0.5
    tf = shape.text_frame
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = 'Helvetica Neue'
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 0, 0)
    return shape


def add_code_block(slide, code, left, top, width, height, *, size=11, bg=CODE_BG, color=TEXT):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg
    rect.line.color.rgb = RGBColor(0x2a, 0x2a, 0x2a)
    rect.line.width = Pt(0.75)
    tf = rect.text_frame
    tf.word_wrap = False
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    lines = code.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.2
        run = p.add_run()
        run.text = line if line else ' '
        run.font.name = 'Menlo'
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return rect


def add_output_block(slide, text, left, top, width, height, *, size=12, color=GOOD):
    """Green-bordered expected-output block."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = OUT_BG
    rect.line.color.rgb = GOOD
    rect.line.width = Pt(2.5)
    tf = rect.text_frame
    tf.word_wrap = False
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line if line else ' '
        run.font.name = 'Menlo'
        run.font.size = Pt(size)
        run.font.color.rgb = RGBColor(0xD8, 0xE8, 0xD8)
    return rect


def add_notes(slide, notes):
    nf = slide.notes_slide.notes_text_frame
    nf.text = notes


def add_image(slide, image_path, left, top, width=None, height=None):
    if width and height:
        return slide.shapes.add_picture(image_path, left, top, width=width, height=height)
    elif width:
        return slide.shapes.add_picture(image_path, left, top, width=width)
    else:
        return slide.shapes.add_picture(image_path, left, top, height=height)


# ============= helper: section header bar =============
def add_section_header(slide, label, *, bg=TAG_BG):
    """Top-left tag and title row."""
    pass


# ============= slide builders =============
def slide_cover():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    add_tag(s, 'WSU · MAY 15, 2026', Inches(0.6), Inches(0.6))
    add_text(s, 'From Black Boxes to Glass Boxes',
             Inches(0.6), Inches(1.5), Inches(12), Inches(1.4),
             size=54, bold=True, color=TEXT, font='Helvetica Neue')
    add_text(s, 'Building Explainable Neural Trees for Safety-Critical Decisions',
             Inches(0.6), Inches(3.0), Inches(12), Inches(0.8),
             size=24, color=ACCENT)
    add_text(s, 'Cagri Temel — CTO, Hezarfen LLC · IEEE Senior Member',
             Inches(0.6), Inches(5.4), Inches(12), Inches(0.5),
             size=16, color=MUTED, bold=True)
    add_text(s, 'Washington State University · Data & Analytics Breakout',
             Inches(0.6), Inches(6.0), Inches(12), Inches(0.4),
             size=12, color=DIM)
    add_notes(s, "Open by introducing yourself in one sentence. "
              "Cagri Temel, CTO of Hezarfen LLC, IEEE Senior Member. "
              "Tell them what they will get in the next 165 minutes: "
              "they will build, break, and explain a real machine learning model "
              "trained on NASA turbofan data. "
              "Set the tone now: this is hands-on, not a passive lecture. "
              "Pause for two seconds after the title slide before advancing.")


def slide_who():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    add_tag(s, 'WHO IS RUNNING THIS', Inches(0.6), Inches(0.6))
    add_text(s, 'Cagri Temel',
             Inches(0.6), Inches(1.3), Inches(10), Inches(0.9),
             size=44, bold=True, color=ACCENT)
    bio = [
        ('CTO of Hezarfen LLC, building explainable AI for regulated industries', {'size': 18}),
        ('IEEE Senior Member · open-source maintainer', {'size': 18}),
        ('Author of neural-trees (PyPI, MIT), the package you will use today', {'size': 18}),
        ('Workshop paper: Robust and Uncertainty-Aware RUL Prediction with Temporal Neural Trees', {'size': 16, 'italic': True, 'color': MUTED}),
        ('Today: in 2 hours and 45 minutes you will build, break, and explain a real model', {'size': 18, 'color': ACCENT, 'bold': True}),
    ]
    add_multi(s, bio, Inches(0.8), Inches(2.5), Inches(12), Inches(4), line_spacing=1.5)
    add_text(s, 'github.com/cgrtml/neural-trees · linkedin.com/in/cagritemel',
             Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
             size=11, color=DIM)
    add_notes(s, "Spend 30 seconds here at most. Bio is on the slide, do not read it line by line. "
              "Emphasize the last bullet: 'In two hours and forty-five minutes, you will build, break, and explain a real model.' "
              "That sets the contract for the day. "
              "Move on as soon as you finish that line.")


def slide_mission():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    add_tag(s, "TODAY'S MISSION", Inches(0.6), Inches(0.6))
    add_text(s, 'Data-driven decisions, when the cost of being wrong is high',
             Inches(0.6), Inches(1.3), Inches(12), Inches(1.6),
             size=34, bold=True, color=ACCENT)
    add_text(s,
             "For the next few hours we will look at how data-driven systems make decisions, "
             "and at how we make those decisions trustworthy in safety-critical settings.",
             Inches(0.6), Inches(3.4), Inches(12), Inches(1.5),
             size=20, color=TEXT)
    add_text(s, "Today's case: a turbofan engine that needs to predict its own failure.",
             Inches(0.6), Inches(5.4), Inches(12), Inches(1),
             size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_notes(s, "This is the framing slide. Say: 'We are going to look at data-driven decisions, "
              "specifically in settings where being wrong is expensive — lives, money, certifications.' "
              "Then announce the case: 'a turbofan engine that needs to predict its own failure.' "
              "Pause briefly before going to the hook slide. The hook slide is intentionally minimal.")


def slide_hook():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    add_text(s, 'A turbofan engine.',
             Inches(0.6), Inches(2.4), Inches(12), Inches(1.2),
             size=56, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, 'When will it fail?',
             Inches(0.6), Inches(3.8), Inches(12), Inches(1.0),
             size=42, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, 'If you knew with certainty, what would you do differently?',
             Inches(0.6), Inches(5.4), Inches(12), Inches(0.6),
             size=18, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    add_notes(s, "Slow this slide down. Read the title aloud: 'A turbofan engine. When will it fail?' "
              "Then stop. Do not fill the silence. Two to three seconds. "
              "Then read the smaller line: 'If you knew with certainty, what would you do differently?' "
              "Let the audience think about it. Do not call on anyone. Move on after a beat.")


def slide_stakes():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    add_tag(s, 'STAKES', Inches(0.6), Inches(0.6))
    add_text(s, 'Getting this wrong is expensive',
             Inches(0.6), Inches(1.3), Inches(12), Inches(1),
             size=36, bold=True, color=TEXT)
    rows = [
        ('Predict too late. ', WARN, 'Engine fails mid-flight. Lives lost.', TEXT),
        ('Predict too early. ', WARN, 'Healthy engines replaced. $100M wasted.', TEXT),
        ('Predict accurately but cannot explain. ', WARN, 'FAA refuses to certify your system.', TEXT),
    ]
    y = Inches(2.6)
    for headline, hc, body, bc in rows:
        tb = slide.shapes.add_textbox(Inches(0.8), y, Inches(12), Inches(0.6)) if False else None
        tf_box = s.shapes.add_textbox(Inches(0.8), y, Inches(12), Inches(0.6))
        tf = tf_box.text_frame
        tf.word_wrap = True
        tf.margin_top = Pt(0)
        tf.margin_bottom = Pt(0)
        p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = headline; r1.font.size = Pt(20); r1.font.bold = True
        r1.font.color.rgb = hc; r1.font.name = 'Helvetica Neue'
        r2 = p.add_run(); r2.text = body; r2.font.size = Pt(20); r2.font.color.rgb = bc
        r2.font.name = 'Helvetica Neue'
        y += Inches(0.75)
    add_text(s, 'Accuracy alone is not enough. You also need to defend the prediction.',
             Inches(0.6), Inches(5.8), Inches(12), Inches(1),
             size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_notes(s, "Three failure modes. Read them aloud, one beat between each. "
              "Notice the third one is the one that surprises engineers: 'accurate but unexplainable still gets rejected.' "
              "Land on the bottom line: 'Accuracy alone is not enough. You also need to defend the prediction.' "
              "This is the thesis of the entire workshop. Say it slowly.")


def slide_cmapss():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    add_tag(s, 'DATASET', Inches(0.6), Inches(0.6))
    add_text(s, 'NASA CMAPSS', Inches(0.6), Inches(1.3), Inches(10), Inches(1),
             size=44, bold=True, color=ACCENT)
    add_text(s, 'Commercial Modular Aero-Propulsion System Simulation',
             Inches(0.6), Inches(2.3), Inches(12), Inches(0.5),
             size=16, italic=True, color=MUTED)
    items = [
        ('100 turbofan engines, simulated to failure', {'size': 18, 'space_after': 8}),
        ('21 sensors per timestep: temperatures, pressures, fan speeds, fuel flow', {'size': 18, 'space_after': 8}),
        ('Each engine logs hundreds of flight cycles before failing', {'size': 18, 'space_after': 8}),
        ('Public, open data, cited in 1,000+ academic papers', {'size': 18, 'space_after': 8}),
    ]
    add_multi(s, items, Inches(0.8), Inches(3.2), Inches(12), Inches(3), line_spacing=1.4)
    add_text(s, 'Source: NASA Prognostics Center of Excellence · ti.arc.nasa.gov',
             Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
             size=11, color=DIM)
    add_notes(s, "Briefly describe CMAPSS: NASA's open dataset, 100 simulated engines, 21 sensors. "
              "Mention that this is the canonical benchmark for predictive maintenance research. "
              "Do not dwell on dataset history. Move on in under a minute. "
              "Students will load and explore this data themselves in Activity 1.")


def slide_rul():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    add_tag(s, 'THE TARGET', Inches(0.6), Inches(0.6))
    add_text(s, 'Remaining Useful Life (RUL)',
             Inches(0.6), Inches(1.3), Inches(12), Inches(1),
             size=38, bold=True, color=ACCENT)
    add_text(s, 'The number of flight cycles left before this engine fails.',
             Inches(0.6), Inches(2.6), Inches(12), Inches(0.6),
             size=20, color=TEXT)
    ascii_art = "cycle 0  →  cycle 50  →  cycle 100  →  cycle 150  →  cycle 192 (failure)\nRUL=192     RUL=142      RUL=92        RUL=42         RUL=0 💥"
    add_code_block(s, ascii_art, Inches(1.0), Inches(3.6), Inches(11.3), Inches(1.6), size=14)
    add_text(s, "Train: complete trajectories. Test: partial trajectories, where you predict RUL from what you can see.",
             Inches(0.6), Inches(5.6), Inches(12), Inches(1),
             size=15, color=MUTED, italic=True)
    add_notes(s, "RUL is the supervised label. For each row in the training set, RUL equals "
              "max_cycle for that engine minus current cycle. Healthy engines have large RUL, "
              "dying engines have small RUL. "
              "We will cap RUL at 125 in the notebook because predicting more than 125 cycles ahead "
              "is not informative. This is standard practice in CMAPSS work.")


def slide_sensors():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    add_tag(s, 'WHAT YOU GET PER CYCLE', Inches(0.6), Inches(0.6))
    add_text(s, '21 sensor channels',
             Inches(0.6), Inches(1.3), Inches(12), Inches(0.9),
             size=36, bold=True, color=ACCENT)
    # Build a simple table-like text grid
    table_rows = [
        ('s2  · T24',  'LPC outlet temperature'),
        ('s3  · T30',  'HPC outlet temperature'),
        ('s4  · T50',  'LPT outlet temperature'),
        ('s7  · P30',  'HPC outlet pressure'),
        ('s8  · Nf',   'Physical fan speed'),
        ('s9  · Nc',   'Physical core speed'),
        ('s11 · Ps30', 'Static pressure HPC outlet'),
        ('s12 · phi',  'Ratio of fuel flow to Ps30'),
        ('s14 · NRc',  'Corrected core speed'),
        ('s15 · BPR',  'Bypass ratio'),
        ('s17 · htBleed', 'Bleed enthalpy'),
        ('s20 · W31',  'HPT coolant bleed'),
    ]
    rows_y = Inches(2.4)
    for i, (left_col, right_col) in enumerate(table_rows):
        row_y = rows_y + Inches(0.32 * i)
        is_highlight = left_col.startswith(('s7', 's11', 's14'))
        col_color = ACCENT if is_highlight else TEXT
        add_text(s, left_col, Inches(2.0), row_y, Inches(3), Inches(0.32), size=14, bold=is_highlight, color=col_color, font='Menlo')
        add_text(s, right_col, Inches(5.5), row_y, Inches(7), Inches(0.32), size=14, color=col_color if is_highlight else MUTED)
    add_text(s, 'Highlighted rows are the most predictive after feature selection.',
             Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
             size=12, color=DIM, italic=True)
    add_notes(s, "Don't read every row. Point at the three highlighted sensors: s7 (HPC outlet pressure), "
              "s11 (static pressure), s14 (corrected core speed). "
              "Tell the audience: 'Remember these three. They will show up again in Activity 1 as the most "
              "important features, and in Activity 2 as the sensors that get attacked.' "
              "Move on in under a minute.")


def slide_lstm():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    add_tag(s, 'A MODERN APPROACH', Inches(0.6), Inches(0.6))
    add_text(s, 'Deep learning gets us close',
             Inches(0.6), Inches(1.3), Inches(12), Inches(1),
             size=38, bold=True, color=ACCENT)
    code = "model = nn.LSTM(input_size=17, hidden_size=64, num_layers=2)\n# … training …\nRMSE on test set: 15.49 cycles\nR²: 0.855"
    add_code_block(s, code, Inches(2.0), Inches(3.0), Inches(9.3), Inches(2.0), size=14)
    add_text(s, "That's pretty good. About 15 cycles off, on average. So we ship it to the FAA, right?",
             Inches(0.6), Inches(5.6), Inches(12), Inches(1),
             size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_notes(s, "Show the result and read it: 'LSTM, two layers, RMSE 15.49 cycles, R-squared 0.855.' "
              "Then deliver the loaded question: 'So we ship it to the FAA, right?' "
              "Lean into the implicit setup. The next slide is the answer.")


def slide_regulator():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s, BG_WARN)
    add_text(s, 'No.', Inches(0.6), Inches(0.6), Inches(12), Inches(1.6),
             size=120, bold=True, color=BAD, align=PP_ALIGN.CENTER)
    add_text(s, "Regulators do not ask 'how accurate is your model?'",
             Inches(0.6), Inches(3.0), Inches(12), Inches(0.8),
             size=24, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, "They ask 'why did your model say that?'",
             Inches(0.6), Inches(3.9), Inches(12), Inches(0.8),
             size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, 'FAA · FDA · OCC (banking) · EU AI Act',
             Inches(0.6), Inches(5.4), Inches(12), Inches(0.5),
             size=18, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, 'All require traceable, explainable decisions for high-risk systems.',
             Inches(0.6), Inches(6.0), Inches(12), Inches(0.5),
             size=14, color=DIM, italic=True, align=PP_ALIGN.CENTER)
    add_notes(s, "Dramatic pause. Just say 'No.' Then count to two in your head before continuing. "
              "'Regulators don't ask how accurate your model is. They ask why your model said that.' "
              "Then list the regulators: FAA for aviation, FDA for medical devices, OCC for banking, EU AI Act for everything in Europe. "
              "Tell them this is not theory. This is what blocks AI systems from shipping in regulated industries.")


def slide_bb_gb():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    add_text(s, 'Black Box vs. Glass Box',
             Inches(0.6), Inches(0.6), Inches(12), Inches(1),
             size=36, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    # Left card
    left_rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.0), Inches(5.5), Inches(3.6))
    left_rect.fill.solid(); left_rect.fill.fore_color.rgb = RGBColor(0x2a, 0x12, 0x12)
    left_rect.line.color.rgb = BAD; left_rect.line.width = Pt(2)
    tf = left_rect.text_frame; tf.margin_left = Pt(20); tf.margin_top = Pt(16)
    items = [('Black Box', 24, True, BAD), ('LSTM, deep CNN, transformer', 14, False, MUTED),
             ('✓ High accuracy', 14, False, GOOD), ('✗ Why? — no answer', 14, False, WARN),
             ('✗ FAA / FDA / OCC: rejected', 14, False, WARN)]
    for i, (text, size, bold, color) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10 if i == 0 else 4)
        r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = 'Helvetica Neue'
    # Right card
    right_rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(2.0), Inches(5.5), Inches(3.6))
    right_rect.fill.solid(); right_rect.fill.fore_color.rgb = RGBColor(0x12, 0x2a, 0x12)
    right_rect.line.color.rgb = GOOD; right_rect.line.width = Pt(2)
    tf = right_rect.text_frame; tf.margin_left = Pt(20); tf.margin_top = Pt(16)
    items = [('Glass Box', 24, True, GOOD), ('Decision tree, linear regression', 14, False, MUTED),
             ('✓ Fully traceable', 14, False, GOOD), ('✗ Lower accuracy', 14, False, WARN),
             ('✓ Regulators: approved', 14, False, GOOD)]
    for i, (text, size, bold, color) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10 if i == 0 else 4)
        r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = 'Helvetica Neue'
    add_text(s, "Today's question: can we have both?",
             Inches(0.6), Inches(6.2), Inches(12), Inches(0.7),
             size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_notes(s, "Frame the binary that students believe in: high accuracy means black box, "
              "interpretability means low accuracy. Walk through the two cards. "
              "End with the question that the next 90 minutes will answer: 'Can we have both?' "
              "Pause briefly. Yes, we can. That's where soft decision trees come in.")


def slide_classic_tree():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    add_tag(s, 'SECTION 2 · CONCEPT', Inches(0.6), Inches(0.6))
    add_text(s, 'Classic Decision Tree',
             Inches(0.6), Inches(1.3), Inches(12), Inches(1),
             size=36, bold=True, color=ACCENT)
    tree = (
        "             [ Sensor 11 (Ps30) > 47.5 ? ]\n"
        "                  /                  \\\n"
        "                 yes                  no\n"
        "      [ Sensor 14 (NRc) > 8138 ? ]  [ HEALTHY ]\n"
        "          /              \\\n"
        "        yes               no\n"
        "       /                   \\\n"
        " [ CRITICAL ]         [ CAUTION ]"
    )
    add_code_block(s, tree, Inches(1.5), Inches(2.5), Inches(10.3), Inches(3.4), size=14)
    add_text(s, 'Beautifully explainable. Often poor accuracy.',
             Inches(0.6), Inches(6.1), Inches(12), Inches(0.5),
             size=18, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, 'Hard splits: a 0.01 difference in a sensor flips the entire path.',
             Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
             size=14, color=WARN, italic=True, align=PP_ALIGN.CENTER)
    add_notes(s, "Classic decision tree, the kind you would draw on a whiteboard. "
              "Show the structure. Say it loud: 'beautifully explainable.' "
              "Then the catch: 'often poor accuracy, because hard splits are brittle.' "
              "Set up the next two slides: hard split versus soft split.")


def slide_hard_split():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    add_text(s, 'Hard split',
             Inches(0.6), Inches(0.8), Inches(12), Inches(1),
             size=44, bold=True, color=ACCENT)
    code = (
        "def hard_split(x, threshold):\n"
        "    if x > threshold:\n"
        "        return 'right'\n"
        "    else:\n"
        "        return 'left'"
    )
    add_code_block(s, code, Inches(2.0), Inches(2.3), Inches(9.3), Inches(2.4), size=16)
    add_text(s, "Output: 'left' or 'right'. Discrete. No gradient.",
             Inches(0.6), Inches(5.2), Inches(12), Inches(0.6),
             size=18, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, 'Cannot train with backpropagation. Tree built greedily, one node at a time.',
             Inches(0.6), Inches(5.9), Inches(12), Inches(0.6),
             size=15, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    add_notes(s, "Students will not type this code. It is here for illustration. "
              "Walk through it: 'if x is bigger than threshold, go right. Otherwise, go left.' "
              "Output is a string. Discrete. No gradient. "
              "Tell them why this matters: 'No gradient means no backpropagation. "
              "You build the tree node by node, greedily, and you cannot train it end to end.'")


def slide_soft_split():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    add_text(s, 'Soft split',
             Inches(0.6), Inches(0.8), Inches(12), Inches(1),
             size=44, bold=True, color=ACCENT)
    code = (
        "def soft_split(x, w, b):\n"
        "    # sigmoid: σ(z) = 1 / (1 + exp(-z))\n"
        "    p_right = sigmoid(w @ x - b)\n"
        "    return p_right       # ∈ [0, 1]"
    )
    add_code_block(s, code, Inches(2.0), Inches(2.3), Inches(9.3), Inches(2.4), size=16)
    add_text(s, 'Output: a probability of going right. Continuous. Gradient flows.',
             Inches(0.6), Inches(5.2), Inches(12), Inches(0.6),
             size=18, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, 'Train the entire tree end-to-end with backprop, like a small neural network.',
             Inches(0.6), Inches(5.9), Inches(12), Inches(0.6),
             size=15, color=GOOD, italic=True, align=PP_ALIGN.CENTER)
    add_notes(s, "Same code, one substitution. The if-else becomes a sigmoid. "
              "Output is now a probability between 0 and 1. "
              "Tell them: 'this is the only change. One line. But because the output is continuous, "
              "gradient flows. You can train the whole tree end-to-end with backpropagation.' "
              "That single change is what makes the rest of the workshop possible.")


def slide_soft_tree_visual():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    add_text(s, 'Soft Decision Tree',
             Inches(0.6), Inches(0.8), Inches(12), Inches(1),
             size=42, bold=True, color=ACCENT)
    tree = (
        "             [ σ(w·x − b) → 0.82 right ]\n"
        "                  /              \\\n"
        "               0.82             0.18\n"
        "                /                  \\\n"
        "   [ σ(w·x − b) → 0.30 ]   [ leaf: P(C, H, W) ]\n"
        "       /        \\\n"
        "     0.30      0.70\n"
        "      /          \\\n"
        "  [ leaf ]    [ leaf ]"
    )
    add_code_block(s, tree, Inches(1.5), Inches(2.3), Inches(10.3), Inches(3.6), size=14)
    add_text(s, 'Every leaf gets a weighted vote based on path probability.',
             Inches(0.6), Inches(6.2), Inches(12), Inches(0.4),
             size=16, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, 'Final prediction = sum over all leaves.',
             Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
             size=14, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    add_notes(s, "Show the structure: every internal node is a sigmoid split. Every leaf has a "
              "class distribution. The final prediction is a weighted sum of all leaves, weighted "
              "by the probability of reaching each leaf along the way. "
              "This is mathematically a smooth function from input to class probabilities, "
              "fully differentiable.")


def slide_best_of_both():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    add_tag(s, 'WHY THIS IS A BIG DEAL', Inches(0.6), Inches(0.6))
    add_text(s, 'You get both',
             Inches(0.6), Inches(1.3), Inches(12), Inches(1),
             size=42, bold=True, color=ACCENT)
    # Left
    left_rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.7), Inches(5.5), Inches(3.2))
    left_rect.fill.solid(); left_rect.fill.fore_color.rgb = RGBColor(0x12, 0x2a, 0x12)
    left_rect.line.color.rgb = GOOD; left_rect.line.width = Pt(2)
    tf = left_rect.text_frame; tf.margin_left = Pt(20); tf.margin_top = Pt(16)
    items = [('From the tree', 22, True, GOOD),
             ('• Path traceability', 14, False, TEXT),
             ('• Per-feature thresholds', 14, False, TEXT),
             ('• Human-readable rules', 14, False, TEXT)]
    for i, (text, size, bold, color) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10 if i == 0 else 6)
        r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = 'Helvetica Neue'
    # Right
    right_rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(2.7), Inches(5.5), Inches(3.2))
    right_rect.fill.solid(); right_rect.fill.fore_color.rgb = RGBColor(0x12, 0x1a, 0x2a)
    right_rect.line.color.rgb = ACCENT; right_rect.line.width = Pt(2)
    tf = right_rect.text_frame; tf.margin_left = Pt(20); tf.margin_top = Pt(16)
    items = [('From the network', 22, True, ACCENT),
             ('• End-to-end training', 14, False, TEXT),
             ('• Composable with other layers', 14, False, TEXT),
             ('• Higher accuracy', 14, False, TEXT)]
    for i, (text, size, bold, color) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10 if i == 0 else 6)
        r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = 'Helvetica Neue'
    add_notes(s, "Summarize the key insight in one breath: 'From the tree, you get path traceability, "
              "per-feature thresholds, human-readable rules. From the network, you get end-to-end training "
              "and the accuracy that comes with it. You get both. That is why soft decision trees matter.'")


def slide_three_properties():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    add_tag(s, 'WHY THIS MATTERS FOR SAFETY-CRITICAL AI', Inches(0.6), Inches(0.6))
    add_text(s, 'Three properties', Inches(0.6), Inches(1.3), Inches(12), Inches(1),
             size=38, bold=True, color=ACCENT)
    items = [
        ('1. Explainability', 22, True, ACCENT, 0),
        ('Every prediction comes with the path that produced it.', 14, False, MUTED, 12),
        ('2. Noise robustness', 22, True, ACCENT, 0),
        ('Soft splits are less brittle to small sensor perturbations than hard splits.', 14, False, MUTED, 12),
        ('3. Sensor failure tolerance', 22, True, ACCENT, 0),
        ('Channel-level dropout during training teaches the model to operate under partial sensor failure.', 14, False, MUTED, 0),
    ]
    tb = s.shapes.add_textbox(Inches(1.0), Inches(2.6), Inches(11.5), Inches(4))
    tf = tb.text_frame; tf.margin_left = Pt(0)
    for i, (text, size, bold, color, after) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(after)
        r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = 'Helvetica Neue'
    add_text(s, 'You will see all three in the next 90 minutes.',
             Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
             size=14, color=DIM, italic=True, align=PP_ALIGN.CENTER)
    add_notes(s, "Three properties matter for safety-critical AI: explainability, noise robustness, "
              "and sensor-failure tolerance. Read each one. Tell the audience: 'You will see all three "
              "demonstrated in the next 90 minutes, not on a slide, on your own laptop.'")


def slide_paper_numbers():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    add_tag(s, 'FROM MY PAPER', Inches(0.6), Inches(0.6))
    add_text(s, 'The numbers', Inches(0.6), Inches(1.3), Inches(12), Inches(1),
             size=42, bold=True, color=ACCENT)
    # Table
    table_data = [
        ['Model', 'Clean RMSE', '30% sensors missing', 'Explainable?'],
        ['LSTM', '15.49', '+89% degradation', '✗'],
        ['Random Forest', '18.07', 'n/a', 'partial'],
        ['Temporal Neural Tree', '15.78', '+17% degradation', '✓'],
    ]
    rows, cols = 4, 4
    tbl_shape = s.shapes.add_table(rows, cols, Inches(1.5), Inches(2.7), Inches(10.3), Inches(2.6))
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(3.0)
    tbl.columns[1].width = Inches(2.0)
    tbl.columns[2].width = Inches(3.3)
    tbl.columns[3].width = Inches(2.0)
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = table_data[r][c]
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = 'Helvetica Neue'
                    run.font.size = Pt(15 if r > 0 else 14)
                    run.font.bold = (r == 0 or r == 3)
                    if r == 0:
                        run.font.color.rgb = ACCENT
                    elif r == 3:
                        run.font.color.rgb = GOOD
                    elif c == 2 and 'degradation' in table_data[r][c] and r == 1:
                        run.font.color.rgb = WARN
                    else:
                        run.font.color.rgb = TEXT
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x1a) if r == 0 else (RGBColor(0x10, 0x20, 0x10) if r == 3 else BG)
    add_text(s, 'Within 0.3 RMSE of LSTM on clean data. 5x more robust under sensor failure.',
             Inches(0.6), Inches(5.8), Inches(12), Inches(0.8),
             size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_notes(s, "Show the numbers. Two key takeaways. "
              "First: on clean data, Temporal Neural Tree is within 0.3 RMSE of an LSTM. Effectively equal. "
              "Second: when 30 percent of sensors are missing, the LSTM degrades by 89 percent. The Temporal Neural Tree degrades by only 17 percent. "
              "Five times more robust. And it is fully explainable. This is the headline result from the paper that motivates the rest of the workshop.")


def slide_package():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    add_tag(s, "THE TOOL YOU'LL USE", Inches(0.6), Inches(0.6))
    add_text(s, 'neural-trees', Inches(0.6), Inches(1.3), Inches(12), Inches(1.2),
             size=56, bold=True, color=ACCENT)
    add_code_block(s, 'pip install neural-trees', Inches(3.0), Inches(3.0), Inches(7.3), Inches(0.7), size=18)
    items = [
        ('• Open source, MIT licensed', {'size': 16, 'space_after': 6}),
        ('• scikit-learn compatible API', {'size': 16, 'space_after': 6}),
        ('• PyTorch backend', {'size': 16, 'space_after': 6}),
        ('• Soft Decision Trees, Hierarchical Mixture of Experts, statistical model comparison', {'size': 16, 'space_after': 6}),
    ]
    add_multi(s, items, Inches(2.0), Inches(4.2), Inches(10), Inches(2.5))
    add_text(s, 'github.com/cgrtml/neural-trees',
             Inches(0.6), Inches(6.8), Inches(12), Inches(0.4),
             size=12, color=DIM, align=PP_ALIGN.CENTER)
    add_notes(s, "Introduce the tool. 'neural-trees, on PyPI. Open source, MIT license. "
              "Scikit-learn compatible API: fit, predict, predict_proba, score. "
              "PyTorch backend. I wrote this library. Today we will all use it together.' "
              "Mention briefly that the GitHub Sprint at the end will let them contribute back.")


# ============= Activity 1 intro + steps =============

def slide_a1_intro():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s, RGBColor(0x0A, 0x1A, 0x2A))
    add_tag(s, 'ACTIVITY 1 · 55 MIN', Inches(0.6), Inches(0.6))
    add_text(s, 'Train Your Own Neural Tree',
             Inches(0.6), Inches(2.0), Inches(12), Inches(1.5),
             size=54, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, 'Open the QR code on the next slide.',
             Inches(0.6), Inches(4.5), Inches(12), Inches(0.7),
             size=22, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, 'No installs. Just Gmail and a browser.',
             Inches(0.6), Inches(5.4), Inches(12), Inches(0.6),
             size=18, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    add_notes(s, "Activity 1 starts now. The tone shifts: you are no longer lecturing, you are running an activity. "
              "Tell them: 'For the next 55 minutes you do the work. I walk around. If anyone gets stuck, raise a hand and I will come.' "
              "Then point to the QR code on the next slide.")


def slide_a1_qr():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s, WHITE)
    add_text(s, 'Scan to start', Inches(0.6), Inches(0.6), Inches(12), Inches(0.8),
             size=36, bold=True, color=BG, align=PP_ALIGN.CENTER)
    qr_path = os.path.join(ROOT, 'assets', 'qr_landing.png')
    if os.path.exists(qr_path):
        add_image(s, qr_path, Inches(5.0), Inches(1.6), width=Inches(3.5))
    add_text(s, 'github.com/cgrtml/wsu-workshop-may15',
             Inches(0.6), Inches(5.4), Inches(12), Inches(0.5),
             size=18, bold=True, color=BG, align=PP_ALIGN.CENTER)
    add_text(s, '1. Scan with phone camera → 2. Tap "Activity 1" → 3. Sign in with Gmail → 4. "Copy to Drive"',
             Inches(0.6), Inches(6.2), Inches(12), Inches(0.5),
             size=14, color=DIM, align=PP_ALIGN.CENTER)
    add_notes(s, "Leave this slide up for two full minutes. Watch the room. "
              "Most students will scan with their phone and open Colab on the laptop using the link. "
              "If anyone looks confused, walk over and help them. "
              "Do not advance until at least 80 percent of the room is on the Colab notebook.")


def slide_a1_overview():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    add_text(s, "Activity 1: what you will do",
             Inches(0.6), Inches(0.7), Inches(12), Inches(0.9),
             size=36, bold=True, color=ACCENT)
    items = [
        ('1. Load & explore CMAPSS FD001 (10 min)', {'size': 17, 'space_after': 8, 'bold': True}),
        ('Plot a few engines. Find which sensors degrade.', {'size': 13, 'space_after': 16, 'color': MUTED, 'italic': True}),
        ('2. Bin RUL into 3 classes (5 min)', {'size': 17, 'space_after': 8, 'bold': True}),
        ('Critical · Caution · Healthy', {'size': 13, 'space_after': 16, 'color': MUTED, 'italic': True}),
        ('3. Train a SoftDecisionTree (15 min)', {'size': 17, 'space_after': 8, 'bold': True}),
        ('depth=4, max_epochs=30, runs on CPU', {'size': 13, 'space_after': 16, 'color': MUTED, 'italic': True}),
        ('4. Evaluate accuracy & confusion matrix (5 min)', {'size': 17, 'space_after': 16, 'bold': True}),
        ('5. Traverse a prediction (15 min)', {'size': 17, 'space_after': 8, 'bold': True, 'color': ACCENT}),
        ('This is the heart of the workshop.', {'size': 13, 'space_after': 0, 'color': ACCENT, 'italic': True}),
    ]
    add_multi(s, items, Inches(0.8), Inches(1.8), Inches(12), Inches(5.5))
    add_notes(s, "Show the agenda for the next 55 minutes. Emphasize step 5 as the heart of the workshop: "
              "they will literally read out the rule the model used. "
              "Most students will get through the first three steps quickly. The traverse-a-prediction step is where the meaning lands. "
              "Promise them that if they walk out remembering anything, it should be that step.")


def add_step_slide(step_num, total, title, code, expected, hint, *, image_path=None, a2=False):
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    tag_text = f'STEP {step_num} / {total}'
    add_tag(s, tag_text, Inches(0.6), Inches(0.5), bg=TAG_BG_A2 if a2 else TAG_BG)
    add_text(s, title, Inches(0.6), Inches(1.05), Inches(12), Inches(0.9),
             size=28, bold=True, color=ACCENT)
    # left: code
    add_text(s, 'PASTE THIS', Inches(0.6), Inches(2.1), Inches(6), Inches(0.4),
             size=11, bold=True, color=ACCENT)
    add_code_block(s, code, Inches(0.6), Inches(2.5), Inches(6.8), Inches(4.4), size=11)
    # right: output or image
    add_text(s, 'YOU SHOULD SEE', Inches(7.7), Inches(2.1), Inches(6), Inches(0.4),
             size=11, bold=True, color=GOOD)
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(7.7), Inches(2.5), width=Inches(5.4)) if False else add_image(s, image_path, Inches(7.7), Inches(2.5), width=Inches(5.4))
    else:
        add_output_block(s, expected, Inches(7.7), Inches(2.5), Inches(5.4), Inches(3.4), size=12)
    if hint:
        add_text(s, hint, Inches(0.6), Inches(7.0), Inches(12), Inches(0.4),
                 size=11, color=MUTED, italic=True)
    return s


# All 11 Activity 1 steps
A1_STEPS = [
    {
        'num': 1, 'title': 'Imports',
        'code': (
            "import numpy as np\n"
            "import pandas as pd\n"
            "import matplotlib.pyplot as plt\n"
            "import seaborn as sns\n"
            "from sklearn.model_selection import train_test_split\n"
            "from sklearn.preprocessing import StandardScaler\n"
            "from sklearn.metrics import classification_report, confusion_matrix\n"
            "from neural_trees import SoftDecisionTree\n"
            "print('All set.')"
        ),
        'expected': 'All set.',
        'hint': 'If you see IndentationError, every line must start at the left edge.',
        'notes': "Step 1, imports. This cell installs nothing on Colab because everything is pre-installed; "
                 "on local Python the conda env should already have neural-trees. "
                 "Say: 'Paste this, run, expect to see All set.' Then walk around. "
                 "Common failure: indentation errors from paste. Fix is Cmd+A then Shift+Tab."
    },
    {
        'num': 2, 'title': 'Load the data',
        'code': (
            "cols = ['engine_id', 'cycle'] + [f'op{i}' for i in range(1, 4)] + [f's{i}' for i in range(1, 22)]\n"
            "train = pd.read_csv('train_FD001.txt', sep=r'\\s+', header=None, names=cols)\n"
            "print(train.shape, '->', train['engine_id'].nunique(), 'engines')\n"
            "train.head()"
        ),
        'expected': "(20631, 26) -> 100 engines\n\nFirst 5 rows of the dataframe\nshown below.",
        'hint': '100 turbofan engines, 20,631 timesteps total.',
        'notes': "Step 2 loads the CMAPSS FD001 file as a pandas dataframe. "
                 "Tell the room: '100 engines, 21 sensors each, 20,631 total rows.' "
                 "If a student gets FileNotFoundError, the previous wget cell did not run."
    },
    {
        'num': 3, 'title': 'Compute Remaining Useful Life',
        'code': (
            "max_cycles = train.groupby('engine_id')['cycle'].max().rename('max_cycle')\n"
            "train = train.merge(max_cycles, on='engine_id')\n"
            "train['RUL'] = (train['max_cycle'] - train['cycle']).clip(upper=125)\n"
            "train[['engine_id','cycle','max_cycle','RUL']].head()"
        ),
        'expected': (
            "   engine_id  cycle  max_cycle  RUL\n"
            "0          1      1        192  125\n"
            "1          1      2        192  125\n"
            "2          1      3        192  125\n"
            "3          1      4        192  125\n"
            "4          1      5        192  125"
        ),
        'hint': 'RUL capped at 125. Engine 1 lived 192 cycles in total.',
        'notes': "Step 3 computes the supervised label. For each row, RUL is the engine's max cycle minus current cycle. "
                 "Capped at 125 because longer predictions are not useful. "
                 "Healthy engines have RUL near 125, dying engines have RUL near 0."
    },
    {
        'num': 4, 'title': "Visualize one engine's degradation",
        'code': (
            "engine = train[train['engine_id'] == 1].set_index('cycle')\n"
            "engine[['s2', 's7', 's11', 's14']].plot(\n"
            "    subplots=True, layout=(2,2), figsize=(12,6))\n"
            "plt.suptitle('Engine 1 sensor degradation', y=1.02)\n"
            "plt.tight_layout(); plt.show()"
        ),
        'expected': '',
        'hint': 's2, s11, s14 trend up. s7 trends down. The engine is wearing.',
        'notes': "Step 4 produces a four-panel plot. Sensor 2, 11, 14 trend up; sensor 7 trends down. "
                 "This is the first aha moment for students: the data does tell the story of wear and tear. "
                 "If a student's plot is missing one of the panels, they likely typed a different engine_id.",
        'image': os.path.join(IMG, 'step4_engine1.png'),
    },
    {
        'num': 5, 'title': 'Bin RUL into three classes',
        'code': (
            "def bin_rul(rul):\n"
            "    if rul < 30: return 0\n"
            "    elif rul < 80: return 1\n"
            "    return 2\n\n"
            "train['health'] = train['RUL'].apply(bin_rul)\n"
            "train['health'].value_counts().sort_index()"
        ),
        'expected': (
            "health\n"
            "0     3000\n"
            "1     5000\n"
            "2    12631\n"
            "Name: count, dtype: int64"
        ),
        'hint': '14% Critical · 24% Caution · 61% Healthy.',
        'notes': "Step 5 turns RUL into a 3-class problem: Critical, Caution, Healthy. "
                 "Tell the audience the class distribution is imbalanced. "
                 "We will use stratify in the train/test split to preserve that ratio."
    },
    {
        'num': 6, 'title': 'Build the feature matrix',
        'code': (
            "all_sensors = [f's{i}' for i in range(1, 22)]\n"
            "informative = [s for s in all_sensors if train[s].std() > 0.001]\n"
            "X = train[informative].values\n"
            "y = train['health'].values\n"
            "scaler = StandardScaler()\n"
            "X = scaler.fit_transform(X)\n"
            "print('X shape:', X.shape, '· y shape:', y.shape)"
        ),
        'expected': 'X shape: (20631, 15) · y shape: (20631,)',
        'hint': '15 of 21 sensors are informative. s1, s5, s10, s16, s18, s19 are constant and dropped.',
        'notes': "Step 6 drops the six constant sensors that carry no information. "
                 "Then standardizes everything to zero mean and unit variance. "
                 "Soft Decision Trees do not strictly require standardization, but the sigmoid splits behave better with it."
    },
    {
        'num': 7, 'title': 'Train / test split',
        'code': (
            "X_train, X_test, y_train, y_test = train_test_split(\n"
            "    X, y,\n"
            "    test_size=0.2,\n"
            "    stratify=y,\n"
            "    random_state=42,\n"
            ")\n"
            "print(X_train.shape, X_test.shape)"
        ),
        'expected': '(16504, 15) (4127, 15)',
        'hint': '80/20 split. stratify=y preserves the class ratio in both sets.',
        'notes': "Step 7, standard train/test split. Stratify by y so the class ratio is preserved. "
                 "Random state 42 so everyone gets the same numbers."
    },
    {
        'num': 8, 'title': 'Train the Soft Decision Tree (30 to 60 sec)',
        'code': (
            "tree = SoftDecisionTree(\n"
            "    depth=4,\n"
            "    max_epochs=30,\n"
            "    learning_rate=0.01,\n"
            "    penalty_coef=1e-3,\n"
            "    verbose=True,\n"
            ")\n"
            "tree.fit(X_train, y_train)\n"
            "acc = (tree.predict(X_test) == y_test).mean()\n"
            "print(f'Test accuracy: {acc:.3f}')"
        ),
        'expected': (
            "Epoch 5/30   loss=0.362  acc=0.851\n"
            "Epoch 10/30  loss=0.355  acc=0.853\n"
            "Epoch 15/30  loss=0.352  acc=0.854\n"
            "Epoch 20/30  loss=0.351  acc=0.855\n"
            "Epoch 25/30  loss=0.350  acc=0.854\n"
            "Epoch 30/30  loss=0.349  acc=0.855\n"
            "Test accuracy: 0.841"
        ),
        'hint': 'Train vs test gap is about 1 point. No overfitting.',
        'notes': "Step 8 is the slow one. 30 to 60 seconds on CPU. While the model trains, "
                 "use the time to tell the audience what is happening: 'PyTorch is doing backprop on 16,504 rows for 30 epochs. "
                 "Once it finishes, we will have a tree that we can read.'"
    },
    {
        'num': 9, 'title': 'Confusion matrix and per-class report',
        'code': (
            "labels = ['Critical', 'Caution', 'Healthy']\n"
            "y_pred = tree.predict(X_test)\n"
            "cm = confusion_matrix(y_test, y_pred)\n"
            "fig, ax = plt.subplots(figsize=(6, 5))\n"
            "sns.heatmap(cm, annot=True, fmt='d', cmap='Blues',\n"
            "            xticklabels=labels, yticklabels=labels, ax=ax)\n"
            "plt.title('Confusion matrix'); plt.show()\n"
            "print(classification_report(y_test, y_pred,\n"
            "      target_names=labels))"
        ),
        'expected': '',
        'hint': '0 false-Criticals on Healthy engines. Only 5 missed-Critical on 600.',
        'notes': "Step 9, the confusion matrix. Two things to highlight on the matrix. "
                 "First, the bottom-left zero: no Healthy engine was wrongly flagged as Critical. "
                 "Second, only 5 of 600 Critical engines were missed and called Healthy. "
                 "In aviation, you want both numbers small. Critical-recall must be high.",
        'image': os.path.join(IMG, 'step9_confusion.png'),
    },
    {
        'num': 10, 'title': 'Which sensor does each node lean on?',
        'code': (
            "splits = tree.get_split_weights()\n"
            "print(f'Internal nodes: {len(splits)}')\n"
            "for i, w in enumerate(splits):\n"
            "    dom = informative[np.argmax(np.abs(w))]\n"
            "    print(f'  Node {i:2d} -> {dom:4s}')"
        ),
        'expected': (
            "Internal nodes: 15\n"
            "  Node  0 -> s7\n"
            "  Node  1 -> s14\n"
            "  Node  2 -> s6\n"
            "  Node  3 -> s14\n"
            "  Node  4 -> s6\n"
            "  Node  5 -> s6\n"
            "  Node  6 -> s12\n"
            "  ..."
        ),
        'hint': 's14 and s6 dominate multiple nodes. These are the model pillars.',
        'notes': "Step 10 is the first peek inside the model. For each of the 15 internal nodes, "
                 "we print the sensor with the largest absolute weight: that node's dominant feature. "
                 "Three sensors will dominate the tree: s7, s14, s6. "
                 "This is feature importance, but computed natively from the model itself, not post-hoc."
    },
    {
        'num': 11, 'title': 'Read one prediction (the heart of the workshop)',
        'code': (
            "idx = 17\n"
            "x = X_test[idx]\n"
            "true = labels[y_test[idx]]\n"
            "pred = labels[tree.predict([x])[0]]\n"
            "proba = tree.predict_proba([x])[0]\n"
            "root = informative[np.argmax(np.abs(splits[0]))]\n\n"
            "print(f'True class:      {true}')\n"
            "print(f'Predicted class: {pred}')\n"
            "print(f'Probabilities:   '\n"
            "      f'Critical={proba[0]:.3f} '\n"
            "      f'Caution={proba[1]:.3f} '\n"
            "      f'Healthy={proba[2]:.3f}')\n"
            "print(f'Root sensor:     {root}')"
        ),
        'expected': (
            "True class:      Caution\n"
            "Predicted class: Healthy\n"
            "Probabilities:   Critical=0.000\n"
            "                 Caution=0.148\n"
            "                 Healthy=0.852\n"
            "Root sensor:     s7"
        ),
        'hint': 'You can defend this prediction to a regulator. An LSTM cannot.',
        'notes': "Step 11 is the punchline of Activity 1. Once everyone has the output, pause and project this slide. "
                 "Walk through it on stage: 'The model predicted Healthy with 85 percent confidence, only 15 percent Caution, "
                 "and zero probability of Critical. It says the root node leaned on sensor 7. "
                 "If a regulator asks why, we can say sensor 7 was above the model's threshold and that drove the decision. "
                 "An LSTM cannot give you this answer. That is the entire point of the workshop.'"
    },
]


def slide_checkpoint():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s, BG_SECTION)
    add_text(s, 'Checkpoint', Inches(0.6), Inches(0.6), Inches(12), Inches(1),
             size=44, bold=True, color=GOOD)
    add_text(s, 'By 10:25 you should have:',
             Inches(0.8), Inches(2.0), Inches(12), Inches(0.6),
             size=20, color=TEXT)
    items = [
        ('✓ A trained SoftDecisionTree on FD001', {'size': 18, 'space_after': 8, 'color': GOOD}),
        ('✓ Test accuracy > 80%', {'size': 18, 'space_after': 8, 'color': GOOD}),
        ('✓ A printed decision path for at least one engine', {'size': 18, 'space_after': 8, 'color': GOOD}),
    ]
    add_multi(s, items, Inches(1.4), Inches(3.0), Inches(12), Inches(2.5))
    add_text(s, '☕ 10-minute break.', Inches(0.6), Inches(5.6), Inches(12), Inches(0.6),
             size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, 'Activity 2 (team competition) starts at 10:35.',
             Inches(0.6), Inches(6.3), Inches(12), Inches(0.5),
             size=14, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    add_notes(s, "Checkpoint. Acknowledge what they have built so far. "
              "Three deliverables: a trained model, accuracy above 80 percent, a decision path printed for one engine. "
              "Announce the break and the start of Activity 2. "
              "During the break, watch for students who did not complete step 11. Help them in the break.")


# ============= Activity 2 =============

def slide_a2_intro():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s, RGBColor(0x1A, 0x0A, 0x1A))
    add_tag(s, 'ACTIVITY 2 · 25 MIN · TEAM CHALLENGE', Inches(0.6), Inches(0.6), bg=TAG_BG_A2)
    add_text(s, 'Sensor Fault Detection',
             Inches(0.6), Inches(2.0), Inches(12), Inches(1.5),
             size=54, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, "One of engine 17's sensors is reporting bad data.",
             Inches(0.6), Inches(4.4), Inches(12), Inches(0.6),
             size=20, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, 'Which sensor, and can your model tell you why?',
             Inches(0.6), Inches(5.2), Inches(12), Inches(0.7),
             size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_notes(s, "Activity 2 is the team competition. Tone shifts again: louder, faster, more energetic. "
              "Set the stakes: 'One sensor is lying. Find which one, and explain how.' "
              "Tell them: '3-4 student teams. 25 minutes. First team to identify all three attacks wins.'")


def slide_a2_scenario():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    add_text(s, 'The scenario',
             Inches(0.6), Inches(0.7), Inches(12), Inches(0.9),
             size=38, bold=True, color=ACCENT)
    add_text(s, 'You are the data science team at an airline. Maintenance ops sends you three suspicious test files for the same engine:',
             Inches(0.6), Inches(1.9), Inches(12), Inches(1.2),
             size=18, color=TEXT)
    quote_rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.2), Inches(11.7), Inches(1.6))
    quote_rect.fill.solid(); quote_rect.fill.fore_color.rgb = RGBColor(0x2A, 0x1F, 0x0A)
    quote_rect.line.color.rgb = WARN; quote_rect.line.width = Pt(2)
    tf = quote_rect.text_frame; tf.margin_left = Pt(16); tf.margin_top = Pt(12); tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = ('Something is wrong with engine 17. The pattern looks like a sensor drift, a frozen sensor, '
              'or noise spikes, but we cannot tell which channel is at fault. Localize the faulty sensor before the next flight.')
    r.font.name = 'Helvetica Neue'; r.font.size = Pt(16); r.font.italic = True; r.font.color.rgb = TEXT
    add_text(s, 'Two models available to you:', Inches(0.6), Inches(5.2), Inches(12), Inches(0.5),
             size=16, bold=True, color=TEXT)
    items = [
        ('• Your SoftDecisionTree from Activity 1, explainable', {'size': 15, 'color': GOOD}),
        ('• A RandomForest baseline, accurate but opaque', {'size': 15, 'color': WARN}),
    ]
    add_multi(s, items, Inches(0.9), Inches(5.8), Inches(12), Inches(1.2))
    add_notes(s, "Read the maintenance ops quote aloud, slowly. Sell the urgency. "
              "Then introduce the two models: their own soft tree from Activity 1, "
              "plus a RandomForest baseline as the black-box comparison. "
              "Tell them the same investigation technique works whether the cause is hardware failure or adversarial attack.")


def slide_a2_teams():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    add_text(s, 'Teams', Inches(0.6), Inches(0.7), Inches(12), Inches(1),
             size=44, bold=True, color=ACCENT)
    items = [
        ('• Form groups of 3 to 4 students', {'size': 20, 'space_after': 12}),
        ('• Use the same notebook within your team', {'size': 20, 'space_after': 12}),
        ('• You have 17 minutes for analysis', {'size': 20, 'space_after': 12}),
        ('• Then we collect answers', {'size': 20, 'space_after': 12}),
    ]
    add_multi(s, items, Inches(1.0), Inches(2.5), Inches(12), Inches(3.5))
    add_text(s, 'First team to correctly identify all three attacks wins.',
             Inches(0.6), Inches(6.0), Inches(12), Inches(0.8),
             size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_notes(s, "Give them 30 seconds to form groups. Most rooms organize themselves; do not micromanage. "
              "Emphasize the time budget (17 minutes of analysis) and the prize (bragging rights for the winning team). "
              "Announce that you will walk around and give hints, but will not give the answer.")


def slide_a2_three_attacks():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    add_text(s, 'Three attack files', Inches(0.6), Inches(0.7), Inches(12), Inches(1),
             size=40, bold=True, color=ACCENT)
    # Table
    rows, cols = 4, 3
    tbl_shape = s.shapes.add_table(rows, cols, Inches(1.0), Inches(2.4), Inches(11.3), Inches(2.6))
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(2.7)
    tbl.columns[1].width = Inches(4.3)
    tbl.columns[2].width = Inches(4.3)
    headers = ['File', 'Attack type', 'Your job']
    rows_data = [
        ['attack_A.csv', 'Drift (constant offset)', 'Which sensor?'],
        ['attack_B.csv', 'Stuck-at (frozen value)', 'Which sensor?'],
        ['attack_C.csv', 'Gaussian noise injection', 'Which sensor?'],
    ]
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c); cell.text = h
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.name = 'Helvetica Neue'; run.font.size = Pt(16); run.font.bold = True
                run.font.color.rgb = ACCENT
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x1a)
    for r, row_data in enumerate(rows_data, 1):
        for c, val in enumerate(row_data):
            cell = tbl.cell(r, c); cell.text = val
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = 'Menlo' if c == 0 else 'Helvetica Neue'
                    run.font.size = Pt(15)
                    run.font.color.rgb = TEXT
            cell.fill.solid(); cell.fill.fore_color.rgb = BG
    add_text(s, 'Each file is a copy of clean engine 17 data, but one sensor channel has been manipulated.',
             Inches(0.6), Inches(5.6), Inches(12), Inches(0.6),
             size=15, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    add_notes(s, "Walk through the three attack files and the corresponding attack types. "
              "Drift: a constant offset added to one sensor for all cycles. "
              "Stuck-at: a sensor frozen at a single value. "
              "Noise: Gaussian noise added to one sensor. "
              "Each file is a corrupted copy of engine 17. Same engine, three different manipulations.")


# Activity 2 step slides
A2_STEPS = [
    {
        'num': 1, 'title': 'Add a RandomForest baseline',
        'code': (
            "from sklearn.ensemble import RandomForestClassifier\n"
            "rf = RandomForestClassifier(\n"
            "    n_estimators=100, random_state=42\n"
            ").fit(X_train, y_train)\n\n"
            "print(f'SoftTree accuracy: '\n"
            "      f'{(tree.predict(X_test) == y_test).mean():.3f}')\n"
            "print(f'RandomForest      : '\n"
            "      f'{(rf.predict(X_test) == y_test).mean():.3f}')"
        ),
        'expected': "SoftTree accuracy: 0.841\nRandomForest      : 0.843",
        'hint': 'Two models, nearly identical accuracy. One is explainable. One is not.',
        'notes': "Step 1, add a RandomForest baseline. The two models reach nearly identical accuracy. "
                 "The difference will appear under attack, not on clean data."
    },
    {
        'num': 2, 'title': 'Predict on clean engine 17',
        'code': (
            "clean = pd.read_csv('engine17_clean.csv')\n"
            "Xc = scaler.transform(clean[informative].values)\n"
            "clean_tree_pred = tree.predict(Xc)\n"
            "clean_rf_pred = rf.predict(Xc)\n\n"
            "print(f'Engine 17 clean: {len(Xc)} cycles')\n"
            "print(f'Tree vs RF agreement: '\n"
            "      f'{(clean_tree_pred == clean_rf_pred).mean():.2f}')"
        ),
        'expected': "Engine 17 clean: 276 cycles\nTree vs RF agreement: 0.95",
        'hint': '95% agreement on clean data. The two models say roughly the same thing.',
        'notes': "Step 2 establishes the clean baseline. The two models agree on 95 percent of cycles. "
                 "We will use this baseline to measure how much each attack changes predictions."
    },
    {
        'num': 3, 'title': 'How much did each attack move the predictions?',
        'code': (
            "for fname in ['attack_A.csv', 'attack_B.csv', 'attack_C.csv']:\n"
            "    df = pd.read_csv(fname)\n"
            "    Xa = scaler.transform(df[informative].values)\n"
            "    pt = tree.predict(Xa)\n"
            "    pr = rf.predict(Xa)\n"
            "    print(f'{fname}: '\n"
            "          f'Tree {(pt != clean_tree_pred).sum():3d}/{len(pt)} | '\n"
            "          f'RF {(pr != clean_rf_pred).sum():3d}/{len(pr)}')"
        ),
        'expected': (
            "attack_A.csv: Tree  18/276 | RF  14/276\n"
            "attack_B.csv: Tree  33/276 | RF   8/276\n"
            "attack_C.csv: Tree   8/276 | RF   1/276"
        ),
        'hint': 'Attack C: RF changed at only 1 cycle. It is missing the attack. The Soft Tree caught 8.',
        'notes': "Step 3 is the first big insight. The RandomForest on Attack C changed its prediction "
                 "at only one cycle out of 276. It is essentially blind to the attack. "
                 "The Soft Tree caught 8 cycles. This is the explainable model not just being interpretable "
                 "but also being more sensitive to subtle manipulations."
    },
    {
        'num': 4, 'title': 'Attack A. Look for two parallel lines',
        'code': (
            "clean = pd.read_csv('engine17_clean.csv')\n"
            "attack = pd.read_csv('attack_A.csv')\n"
            "fig, axes = plt.subplots(4, 4, figsize=(16, 10))\n"
            "for ax, s in zip(axes.flat, informative):\n"
            "    ax.plot(clean['cycle'], clean[s],\n"
            "            label='clean', alpha=0.7)\n"
            "    ax.plot(attack['cycle'], attack[s],\n"
            "            label='attack', alpha=0.7)\n"
            "    ax.set_title(s)\n"
            "axes.flat[0].legend()\n"
            "plt.tight_layout(); plt.show()"
        ),
        'expected': '',
        'hint': 'One sensor shows the orange line shifted up by a constant offset. Answer: s11 (drift).',
        'notes': "Step 4 visualizes attack A. Drift looks like two parallel lines: the attack copy of one sensor "
                 "is shifted up by a constant amount. Answer is sensor 11.",
        'image': os.path.join(IMG, 'attack_A.png'),
    },
    {
        'num': 5, 'title': 'Attack B. Look for a flat horizontal line',
        'code': (
            "attack = pd.read_csv('attack_B.csv')\n"
            "fig, axes = plt.subplots(4, 4, figsize=(16, 10))\n"
            "for ax, s in zip(axes.flat, informative):\n"
            "    ax.plot(clean['cycle'], clean[s],\n"
            "            label='clean', alpha=0.7)\n"
            "    ax.plot(attack['cycle'], attack[s],\n"
            "            label='attack', alpha=0.7)\n"
            "    ax.set_title(s)\n"
            "axes.flat[0].legend()\n"
            "plt.tight_layout(); plt.show()"
        ),
        'expected': '',
        'hint': "One sensor's orange line is flat while blue rises. Stuck-at. Answer: s14 (stuck-at).",
        'notes': "Step 5 visualizes attack B. Stuck-at: one sensor reports the same value for every cycle "
                 "while the real value would be evolving. Answer is sensor 14, which is one of the model's "
                 "pillars (we saw it in Activity 1, Step 10). Attackers go for what the model relies on.",
        'image': os.path.join(IMG, 'attack_B.png'),
    },
    {
        'num': 6, 'title': 'Attack C. Look for a noisier line',
        'code': (
            "attack = pd.read_csv('attack_C.csv')\n"
            "fig, axes = plt.subplots(4, 4, figsize=(16, 10))\n"
            "for ax, s in zip(axes.flat, informative):\n"
            "    ax.plot(clean['cycle'], clean[s],\n"
            "            label='clean', alpha=0.7)\n"
            "    ax.plot(attack['cycle'], attack[s],\n"
            "            label='attack', alpha=0.7)\n"
            "    ax.set_title(s)\n"
            "axes.flat[0].legend()\n"
            "plt.tight_layout(); plt.show()"
        ),
        'expected': '',
        'hint': 'Same trend, much higher variance. Noise injection. Answer: s9 (noise).',
        'notes': "Step 6 visualizes attack C. Noise injection is the subtlest of the three: the trend is preserved, "
                 "only the variance increases. Hardest to detect visually. Answer is sensor 9.",
        'image': os.path.join(IMG, 'attack_C.png'),
    },
    {
        'num': 7, 'title': 'Quantify the divergence',
        'code': (
            "for fname in ['attack_A.csv', 'attack_B.csv', 'attack_C.csv']:\n"
            "    a = pd.read_csv(fname)\n"
            "    d = {s: float(np.mean(np.abs(clean[s].values - a[s].values)))\n"
            "         for s in informative}\n"
            "    top = sorted(d.items(), key=lambda kv: -kv[1])[:3]\n"
            "    print(f'{fname} top 3: {top}')"
        ),
        'expected': (
            "attack_A.csv top 3: [('s11', 0.367), ('s2', 0.0), ('s3', 0.0)]\n"
            "attack_B.csv top 3: [('s14', 19.89), ('s2', 0.0), ('s3', 0.0)]\n"
            "attack_C.csv top 3: [('s9',  6.42), ('s2', 0.0), ('s3', 0.0)]"
        ),
        'hint': 'Numbers confirm the visual answer. Only one sensor differs in each file.',
        'notes': "Step 7 quantifies what they already saw. Each attack file differs from clean engine 17 "
                 "in exactly one sensor: 11, 14, or 9. The other sensors have zero divergence because they "
                 "were not touched. Numbers are unambiguous."
    },
]


def slide_a2_clock():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s, RGBColor(0x10, 0x0A, 0x1A))
    add_text(s, '⏱  17 minutes',
             Inches(0.6), Inches(2.0), Inches(12), Inches(1.6),
             size=80, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, 'Open activity2_student.ipynb',
             Inches(0.6), Inches(4.2), Inches(12), Inches(0.7),
             size=22, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, 'Run cells. Discuss as a team. Pick your answers.',
             Inches(0.6), Inches(5.0), Inches(12), Inches(0.6),
             size=18, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    add_text(s, "I'll come around with hints, but I will not give you the answer.",
             Inches(0.6), Inches(6.2), Inches(12), Inches(0.4),
             size=14, color=DIM, italic=True, align=PP_ALIGN.CENTER)
    add_notes(s, "Start a 17-minute timer. Visible on stage if possible. "
              "Walk around the room. Provoke wrong answers: when a team says 'sensor 11', "
              "ask 'are you sure? attack B has sensor 14 with score 19.89, isn't that bigger?' "
              "Force them to explain their reasoning back to you.")


def slide_a2_answers():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    add_text(s, 'Answers', Inches(0.6), Inches(0.6), Inches(12), Inches(1),
             size=44, bold=True, color=ACCENT)
    rows, cols = 4, 3
    tbl_shape = s.shapes.add_table(rows, cols, Inches(2.0), Inches(2.2), Inches(9.3), Inches(2.8))
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(2.0)
    tbl.columns[1].width = Inches(4.0)
    tbl.columns[2].width = Inches(3.3)
    headers = ['Attack', 'Sensor', 'Type']
    rows_data = [
        ['A', 'Sensor 11 (Ps30)', 'Drift'],
        ['B', 'Sensor 14 (NRc)', 'Stuck-at'],
        ['C', 'Sensor 9 (Nc)', 'Gaussian noise'],
    ]
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c); cell.text = h
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.name = 'Helvetica Neue'; run.font.size = Pt(18); run.font.bold = True
                run.font.color.rgb = ACCENT
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x1a)
    for r, row_data in enumerate(rows_data, 1):
        for c, val in enumerate(row_data):
            cell = tbl.cell(r, c); cell.text = val
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = 'Helvetica Neue'; run.font.size = Pt(20); run.font.color.rgb = TEXT
            cell.fill.solid(); cell.fill.fore_color.rgb = BG
    add_text(s, 'Which team got all three? 🏆',
             Inches(0.6), Inches(5.6), Inches(12), Inches(0.8),
             size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_notes(s, "Reveal the answers. Read them out: sensor 11 drift, sensor 14 stuck-at, sensor 9 noise. "
              "Ask who got all three. Applaud the winning team. Energy spike, then we transition into the analysis slide.")


def slide_what_just_happened():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    add_text(s, 'What just happened',
             Inches(0.6), Inches(0.7), Inches(12), Inches(1),
             size=40, bold=True, color=ACCENT)
    items = [
        ('The RandomForest changed its prediction, but could not tell you why.', {'size': 18, 'space_after': 16}),
        ("The SoftDecisionTree's split weights shifted in a sensor-specific way.", {'size': 18, 'space_after': 16}),
        ('That shift is the explanation. It is a fingerprint of which sensor the model is leaning on.', {'size': 18, 'space_after': 16}),
    ]
    add_multi(s, items, Inches(1.0), Inches(2.2), Inches(12), Inches(3.5))
    add_text(s, 'Explainability did not just make the model auditable.',
             Inches(0.6), Inches(5.7), Inches(12), Inches(0.6),
             size=20, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, 'It made the model a debugging tool for the system around it.',
             Inches(0.6), Inches(6.3), Inches(12), Inches(0.6),
             size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_notes(s, "The teaching point. Explainability is not just about audits. It is also a debugging tool. "
              "If a sensor goes bad in production, the model itself can tell you which sensor. "
              "The same technique that satisfies the FAA also helps your maintenance engineers find the broken hardware.")


def slide_sprint_intro():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s, BG_SECTION)
    add_tag(s, 'SPRINT · 15 MIN', Inches(0.6), Inches(0.6), bg=GOOD)
    add_text(s, '🔧 GitHub Contribution Sprint',
             Inches(0.6), Inches(2.0), Inches(12), Inches(1.5),
             size=48, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, 'Walk out of this room with a real open-source contribution',
             Inches(0.6), Inches(4.3), Inches(12), Inches(0.6),
             size=22, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, 'on your GitHub profile.',
             Inches(0.6), Inches(4.9), Inches(12), Inches(0.6),
             size=22, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, '15 minutes. Your name on a published ML library.',
             Inches(0.6), Inches(6.0), Inches(12), Inches(0.5),
             size=16, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    add_notes(s, "Activity 3 is a different kind of value. Tell them: 'For the last 15 minutes, you are going to make a "
              "real open-source contribution. Not a workshop exercise. A real pull request to a real library that is on PyPI. "
              "Your name on it, on GitHub, forever.' "
              "Energy should be high here. This is also the segment that builds your EB-1A evidence.")


def slide_sprint_mechanics():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    add_text(s, 'How it works', Inches(0.6), Inches(0.7), Inches(12), Inches(1),
             size=40, bold=True, color=ACCENT)
    items = [
        ('1. Open github.com/cgrtml/neural-trees → Issues tab', {'size': 18, 'space_after': 14}),
        ('2. Filter by label: good first issue. Pick one (10 to 20 min scope)', {'size': 18, 'space_after': 14}),
        ('3. Comment "I am taking this" so no duplicates', {'size': 18, 'space_after': 14}),
        ('4. Fork the repo, edit, open a PR', {'size': 18, 'space_after': 14}),
        ('5. I will review live. Simple PRs get merged on the spot.', {'size': 18, 'space_after': 14, 'color': GOOD}),
    ]
    add_multi(s, items, Inches(0.8), Inches(2.0), Inches(12), Inches(5))
    add_text(s, '~20 issues waiting · pick what fits your level',
             Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
             size=14, color=ACCENT, italic=True, align=PP_ALIGN.CENTER)
    add_notes(s, "Walk through the mechanics. Five steps. Emphasize step 5: 'I will review live. "
              "If your PR is correct, I merge it on the spot, and your name goes into the contributors list of a real ML library.' "
              "Roughly 20 good first issues are waiting; they range from doc fixes to small tests to runnable examples.")


def slide_sprint_why():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    add_text(s, 'Why this is worth your 15 minutes',
             Inches(0.6), Inches(0.7), Inches(12), Inches(1),
             size=34, bold=True, color=ACCENT)
    items = [
        ('A concrete open-source contribution on your profile, visible to recruiters, internships, and grad-school admissions.', {'size': 17, 'space_after': 14}),
        ('You become part of a published Python ML library (PyPI: neural-trees).', {'size': 17, 'space_after': 14}),
        ('Your name goes in the README. I will add a "WSU Workshop Contributors" section listing every merged PR.', {'size': 17, 'space_after': 14}),
        ('It is real proof that you can collaborate on a codebase, not just clone tutorials.', {'size': 17, 'space_after': 14}),
    ]
    add_multi(s, items, Inches(0.8), Inches(2.0), Inches(12), Inches(4.5))
    add_text(s, 'Open the repo. Pick an issue. Go.',
             Inches(0.6), Inches(6.4), Inches(12), Inches(0.6),
             size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_notes(s, "Sell the value: a real OSS contribution is a credential that opens doors. "
              "Recruiters see it on GitHub. Internship applications. Grad school. "
              "Every merged PR also gets credited in the neural-trees README under WSU Workshop Contributors. "
              "Then say: 'Open the repo. Pick an issue. Go.' Then walk to your laptop and start reviewing.")


def slide_eu_ai_act():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    add_tag(s, 'REGULATORY CONTEXT', Inches(0.6), Inches(0.6), bg=WARN)
    add_text(s, 'The EU AI Act',
             Inches(0.6), Inches(1.3), Inches(12), Inches(1),
             size=42, bold=True, color=ACCENT)
    add_text(s, 'Took effect in 2025. Applies to any AI system serving EU customers.',
             Inches(0.6), Inches(2.5), Inches(12), Inches(0.6),
             size=18, color=TEXT)
    add_text(s, 'For "high-risk" systems (predictive maintenance for aviation qualifies):',
             Inches(0.6), Inches(3.4), Inches(12), Inches(0.6),
             size=18, color=TEXT)
    items = [
        ('Article 13: Transparency & explanation requirements', {'size': 17, 'space_after': 8, 'bold': True}),
        ('Article 14: Human oversight with meaningful interpretation', {'size': 17, 'space_after': 8, 'bold': True}),
        ('Article 15: Robustness against adversarial inputs and faults', {'size': 17, 'space_after': 8, 'bold': True}),
    ]
    add_multi(s, items, Inches(1.0), Inches(4.4), Inches(12), Inches(2))
    add_text(s, 'Same direction in the US: SR 11-7 (banking), FDA SaMD, FAA AC 25-1309.',
             Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
             size=14, color=DIM, italic=True, align=PP_ALIGN.CENTER)
    add_notes(s, "Connect the workshop to the regulatory landscape. EU AI Act took effect in 2025. "
              "Three articles matter for what we did today: Article 13 is transparency, "
              "Article 14 is human oversight, Article 15 is robustness. "
              "Tell them this is happening in the US too: SR 11-7 for banking, FDA for medical devices, FAA for aviation. "
              "Same direction, different agencies.")


def slide_wrap_up():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s, BG_SECTION)
    add_tag(s, 'WRAP-UP · 5 MIN', Inches(0.6), Inches(0.6), bg=GOOD)
    add_text(s, 'What you walk out with',
             Inches(0.6), Inches(1.3), Inches(12), Inches(1),
             size=38, bold=True, color=TEXT)
    items = [
        ('1. Explainability is mandatory in safety-critical AI.', {'size': 20, 'bold': True, 'space_after': 4, 'color': ACCENT}),
        ('Required by EU AI Act, FDA, FAA, OCC for high-risk systems.', {'size': 14, 'color': MUTED, 'italic': True, 'space_after': 18}),
        ('2. Soft Decision Trees give you both worlds.', {'size': 20, 'bold': True, 'space_after': 4, 'color': ACCENT}),
        ('Neural-network accuracy with tree-level interpretability.', {'size': 14, 'color': MUTED, 'italic': True, 'space_after': 18}),
        ('3. You ran this yourself, on real NASA data.', {'size': 20, 'bold': True, 'space_after': 4, 'color': ACCENT}),
        ('Trained a model · localized a sensor fault · maybe shipped an OSS contribution.', {'size': 14, 'color': MUTED, 'italic': True}),
    ]
    add_multi(s, items, Inches(1.0), Inches(2.5), Inches(12), Inches(4.5))
    add_notes(s, "Three takeaways. Read each headline aloud. The subline explains it but does not need to be read verbatim. "
              "Lean on point 3: 'You ran this. Not me. Not a video. You.' "
              "Then transition into the career bridge.")


def slide_career_bridge():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    add_text(s, 'Where this is going',
             Inches(0.6), Inches(0.7), Inches(12), Inches(1),
             size=40, bold=True, color=ACCENT)
    add_text(s, 'Banks (Model Risk Management), insurers, healthcare, aerospace.',
             Inches(0.6), Inches(2.0), Inches(12), Inches(0.7),
             size=20, color=TEXT)
    add_text(s, 'Every regulated industry needs explainable AI.',
             Inches(0.6), Inches(2.8), Inches(12), Inches(0.7),
             size=20, color=TEXT)
    add_text(s, 'My company Hezarfen builds compliance-ready AI for US mid-market banks under SR 11-7 and the EU AI Act.',
             Inches(0.6), Inches(4.0), Inches(12), Inches(1.5),
             size=18, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, 'If any of this interested you, whether as a thesis topic, an open-source contribution, or just a conversation:',
             Inches(0.6), Inches(5.8), Inches(12), Inches(0.8),
             size=16, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    add_text(s, 'reach out.',
             Inches(0.6), Inches(6.6), Inches(12), Inches(0.6),
             size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_notes(s, "Open the door. Tell them where this work goes next: banking model risk, insurance, healthcare, aerospace. "
              "Mention Hezarfen and the SR 11-7 angle. "
              "Then make the explicit invitation: 'Thesis topic, open-source contribution, conversation. Reach out.' "
              "This is your funnel-top moment for the strongest students in the room.")


def slide_links():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    add_text(s, 'Stay in touch',
             Inches(0.6), Inches(0.7), Inches(12), Inches(1),
             size=40, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, 'Cagri Temel',
             Inches(0.6), Inches(2.4), Inches(12), Inches(0.8),
             size=32, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, 'Hezarfen LLC · IEEE Senior Member',
             Inches(0.6), Inches(3.2), Inches(12), Inches(0.5),
             size=16, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, 'github.com/cgrtml/neural-trees',
             Inches(0.6), Inches(4.3), Inches(12), Inches(0.6),
             size=22, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, 'linkedin.com/in/cagritemel',
             Inches(0.6), Inches(5.1), Inches(12), Inches(0.6),
             size=22, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, 'Workshop materials: github.com/cgrtml/wsu-workshop-may15',
             Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
             size=13, color=DIM, italic=True, align=PP_ALIGN.CENTER)
    add_notes(s, "Leave this slide up while you take questions. Students will photograph it. "
              "Read the three links out loud once.")


def slide_thank_you():
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    add_text(s, 'Thank you',
             Inches(0.6), Inches(2.0), Inches(12), Inches(1.6),
             size=80, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, 'To Dr. Sergey Lapin and Jeremy for the invitation.',
             Inches(0.6), Inches(4.4), Inches(12), Inches(0.6),
             size=22, color=ACCENT, italic=True, align=PP_ALIGN.CENTER)
    add_text(s, 'Questions? I will stay 5 minutes after.',
             Inches(0.6), Inches(6.0), Inches(12), Inches(0.5),
             size=16, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    add_notes(s, "Close with the thank-you. Acknowledge Sergey and Jeremy by name. "
              "Offer five minutes for questions after the session officially ends. "
              "Stay near the stage; students will come up.")


# ============= build the deck =============
def add_a1_step(s_data):
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    tag_text = f"STEP {s_data['num']} / 11"
    add_tag(s, tag_text, Inches(0.6), Inches(0.45))
    add_text(s, s_data['title'], Inches(0.6), Inches(1.0), Inches(12), Inches(0.9),
             size=26, bold=True, color=ACCENT)
    add_text(s, 'PASTE THIS', Inches(0.6), Inches(2.05), Inches(6), Inches(0.4),
             size=11, bold=True, color=ACCENT)
    add_code_block(s, s_data['code'], Inches(0.6), Inches(2.45), Inches(6.8), Inches(4.5), size=11)
    add_text(s, 'YOU SHOULD SEE', Inches(7.7), Inches(2.05), Inches(6), Inches(0.4),
             size=11, bold=True, color=GOOD)
    if 'image' in s_data and os.path.exists(s_data['image']):
        add_image(s, s_data['image'], Inches(7.7), Inches(2.45), width=Inches(5.4))
    else:
        add_output_block(s, s_data['expected'], Inches(7.7), Inches(2.45), Inches(5.4), Inches(3.4), size=12)
    if s_data.get('hint'):
        add_text(s, s_data['hint'], Inches(0.6), Inches(7.05), Inches(12), Inches(0.35),
                 size=11, color=MUTED, italic=True)
    add_notes(s, s_data['notes'])


def add_a2_step(s_data):
    s = prs.slides.add_slide(prs.slide_layouts[BLANK])
    set_bg(s)
    tag_text = f"STEP {s_data['num']} / 7"
    add_tag(s, tag_text, Inches(0.6), Inches(0.45), bg=TAG_BG_A2)
    add_text(s, s_data['title'], Inches(0.6), Inches(1.0), Inches(12), Inches(0.9),
             size=26, bold=True, color=ACCENT)
    add_text(s, 'PASTE THIS', Inches(0.6), Inches(2.05), Inches(6), Inches(0.4),
             size=11, bold=True, color=ACCENT)
    add_code_block(s, s_data['code'], Inches(0.6), Inches(2.45), Inches(6.8), Inches(4.5), size=11)
    add_text(s, 'YOU SHOULD SEE', Inches(7.7), Inches(2.05), Inches(6), Inches(0.4),
             size=11, bold=True, color=GOOD)
    if 'image' in s_data and os.path.exists(s_data['image']):
        add_image(s, s_data['image'], Inches(7.7), Inches(2.45), width=Inches(5.4))
    else:
        add_output_block(s, s_data['expected'], Inches(7.7), Inches(2.45), Inches(5.4), Inches(3.4), size=12)
    if s_data.get('hint'):
        add_text(s, s_data['hint'], Inches(0.6), Inches(7.05), Inches(12), Inches(0.35),
                 size=11, color=MUTED, italic=True)
    add_notes(s, s_data['notes'])


# === build order ===
slide_cover()
slide_who()
slide_mission()
slide_hook()
slide_stakes()
slide_cmapss()
slide_rul()
slide_sensors()
slide_lstm()
slide_regulator()
slide_bb_gb()
slide_classic_tree()
slide_hard_split()
slide_soft_split()
slide_soft_tree_visual()
slide_best_of_both()
slide_three_properties()
slide_paper_numbers()
slide_package()
slide_a1_intro()
slide_a1_qr()
slide_a1_overview()
for step in A1_STEPS:
    add_a1_step(step)
slide_checkpoint()
slide_a2_intro()
slide_a2_scenario()
slide_a2_teams()
slide_a2_three_attacks()
for step in A2_STEPS:
    add_a2_step(step)
slide_a2_clock()
slide_a2_answers()
slide_what_just_happened()
slide_sprint_intro()
slide_sprint_mechanics()
slide_sprint_why()
slide_eu_ai_act()
slide_wrap_up()
slide_career_bridge()
slide_links()
slide_thank_you()

prs.save(OUT_PATH)
print(f'Saved: {OUT_PATH}')
print(f'Total slides: {len(prs.slides)}')
